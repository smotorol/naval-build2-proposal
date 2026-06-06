# -*- coding: utf-8 -*-
from __future__ import annotations

import math
from pathlib import Path
from typing import Iterable, List, Optional, Sequence, Tuple

import fitz
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_CONNECTOR, MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, MSO_AUTO_SIZE, PP_ALIGN
from pptx.util import Cm, Pt


ROOT = Path(__file__).resolve().parent
FORM_PPTX = ROOT / "form" / "정성제안서.pptx"
OCR_RFP_PDF = ROOT / "해군교전2요청서_word.pdf"
RFP_PDF = OCR_RFP_PDF if OCR_RFP_PDF.exists() else ROOT / "해군교전2요청서.pdf"
TASK_TXT = ROOT / "task.txt"
BASE_PPTX = ROOT / "output" / "해군교전분석모델_BuildII_결과분석_제안서_초안.pptx"
OUT = ROOT / "해군교전2_결과분석_제안서_고도화.pptx"
ASSETS = ROOT / "output" / "assets"

NAVY = RGBColor(0x00, 0x4B, 0x73)
BLUE = RGBColor(0x00, 0x79, 0xA8)
SKY = RGBColor(0xD8, 0xEB, 0xF4)
PALE = RGBColor(0xF4, 0xF8, 0xFA)
GRAY = RGBColor(0x66, 0x66, 0x66)
DARK = RGBColor(0x23, 0x2A, 0x31)
LINE = RGBColor(0xC8, 0xD2, 0xDA)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
GREEN = RGBColor(0x69, 0x9B, 0x7A)
RED = RGBColor(0xB8, 0x5C, 0x5C)


def ensure_assets() -> None:
    ASSETS.mkdir(parents=True, exist_ok=True)


def read_task() -> str:
    raw = TASK_TXT.read_bytes()
    for enc in ("utf-8", "cp949", "euc-kr"):
        try:
            return raw.decode(enc)
        except UnicodeDecodeError:
            pass
    return raw.decode("utf-8", errors="replace")


def extract_rfp_references() -> dict[str, List[str]]:
    """Use the OCR RFP text as proposal evidence without rendering PDF pages."""
    ensure_assets()
    doc = fitz.open(RFP_PDF)
    page_text: dict[int, str] = {}
    for page_no in [54, 56, 57, 59, 60, 61, 62, 63]:
        if page_no <= doc.page_count:
            text = doc[page_no - 1].get_text("text")
            cleaned = " ".join(line.strip() for line in text.splitlines() if line.strip())
            page_text[page_no] = cleaned

    refs = {
        "understanding": [
            "p.63 결과분석 데이터 3D 가시화 개선",
            "p.63 자동 3D 재생, 멀티뷰, 계산모델 가시화 요구",
            "p.63 탐지확률·명중확률 등 수학적 계산 결과 실시간 차트 표현",
        ],
        "arch": [
            "p.62 Web 기반 전장상황 전시 및 3차원 기동/탐지/교전 상황 전시",
            "p.57 플랫폼 3D 및 시스템 네트워크 모듈화·자동생성 요구",
        ],
        "3d": [
            "p.63 전장 환경과 플랫폼 및 무기체계 상호작용 전시",
            "p.62 다수 플랫폼 동시 가시화 및 이벤트 기반 과업 진행 전시",
        ],
        "underwater": [
            "p.63 고해상도 수중환경 3D 가시화",
            "p.63 수온 분포, 음파전달손실, 음선경로 반영",
        ],
        "sound": [
            "p.63 음파전달손실 및 음선경로 3D 가시화",
            "p.59 Multi-Static Sonar 다중 음향경로 반영 요구",
        ],
        "model": [
            "p.63 계산모델 가시화",
            "p.56 Look-up Table 데이터 가시화 및 파라미터 자동조정 요구",
        ],
        "detect": [
            "p.63 탐지확률 실시간 차트 표현",
            "p.59 클러터·재밍 효과 반영 탐지확률 계산모델",
            "p.63 경비구역 탐지확률 히트맵 등 가시화",
        ],
        "hit": [
            "p.63 명중확률 실시간 차트 표현",
            "p.61 유도탄 탐지·추적·회피·교전 성공확률 계산",
        ],
        "missile": [
            "p.63 유도탄·기만기 기만 성공확률 산출 개선",
            "p.60 유도탄 기만기 다수·다종 결합 RCS 변화 반영",
        ],
        "torpedo": [
            "p.63 어뢰·기만기 기만 성공확률 산출 개선",
            "p.60 어뢰 기만기 운용모드 및 기만 논리 반영",
        ],
        "ew": [
            "p.63 전자전 장비 효과 산출 추가",
            "p.63 재머, GPS 교란 등 전자전 효과 반영",
        ],
        "flow": [
            "p.63 결과 데이터 HDF 형태 추출 기능 제공",
            "p.62 이벤트에 따른 과업 진행 전시",
        ],
    }

    md = [
        "# 요청서 요구사항 근거 정리",
        "",
        f"- Source: `{RFP_PDF.name}`",
        "- Note: PDF page rendering images are not used. References below are extracted from OCR text and normalized for proposal wording.",
        "",
    ]
    for key, items in refs.items():
        md.append(f"## {key}")
        md.extend(f"- {item}" for item in items)
        md.append("")
    md.append("## Extracted OCR Page Text")
    for page_no, text in page_text.items():
        md.append(f"### p.{page_no}")
        md.append(text[:2200])
        md.append("")
    (ROOT / "output" / "extracted_rfp_text_refs.md").write_text("\n".join(md), encoding="utf-8")
    return refs


def slide_size() -> Tuple[int, int]:
    if FORM_PPTX.exists():
        prs = Presentation(FORM_PPTX)
        return prs.slide_width, prs.slide_height
    return Cm(19.05), Cm(27.5)


def set_text(shape, text: str, size=9, bold=False, color=DARK, align=None) -> None:
    tf = shape.text_frame
    tf.clear()
    tf.word_wrap = True
    tf.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
    tf.margin_left = Cm(0.12)
    tf.margin_right = Cm(0.12)
    tf.margin_top = Cm(0.06)
    tf.margin_bottom = Cm(0.06)
    p = tf.paragraphs[0]
    p.text = text
    p.font.name = "Malgun Gothic"
    p.font.size = Pt(size)
    p.font.bold = bold
    p.font.color.rgb = color
    p.space_after = Pt(0)
    if align:
        p.alignment = align


def textbox(slide, x, y, w, h, text, size=9, bold=False, color=DARK, align=None):
    sh = slide.shapes.add_textbox(int(x), int(y), int(w), int(h))
    set_text(sh, text, size, bold, color, align)
    return sh


def box(slide, x, y, w, h, text="", fill=WHITE, line=LINE, size=8, bold=False, color=DARK, align=PP_ALIGN.CENTER):
    sh = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, int(x), int(y), int(w), int(h))
    sh.fill.solid()
    sh.fill.fore_color.rgb = fill
    sh.line.color.rgb = line
    set_text(sh, text, size, bold, color, align)
    sh.text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
    return sh


def rect(slide, x, y, w, h, fill, line=None):
    sh = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, int(x), int(y), int(w), int(h))
    sh.fill.solid()
    sh.fill.fore_color.rgb = fill
    sh.line.color.rgb = line or fill
    return sh


def line(slide, x1, y1, x2, y2, color=BLUE, width=1.5):
    c = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, int(x1), int(y1), int(x2), int(y2))
    c.line.color.rgb = color
    c.line.width = Pt(width)
    return c


def arrow(slide, x, y, w, h, color=BLUE):
    sh = slide.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, int(x), int(y), int(w), int(h))
    sh.fill.solid()
    sh.fill.fore_color.rgb = color
    sh.line.color.rgb = color
    return sh


def add_header(slide, num: int, title: str, message: str, width, height):
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = WHITE
    textbox(slide, Cm(0.7), Cm(0.55), Cm(2.1), Cm(0.35), f"Ⅱ - {num-1}", 7.5, True, NAVY)
    textbox(slide, Cm(0.7), Cm(0.95), width - Cm(1.4), Cm(0.62), title, 12.5, True, DARK)
    rect(slide, Cm(0.7), Cm(1.78), width - Cm(1.4), Cm(0.03), NAVY)
    rect(slide, Cm(0.7), Cm(1.84), Cm(4.3), Cm(0.07), BLUE)
    textbox(slide, Cm(0.9), Cm(2.12), width - Cm(1.8), Cm(0.56), message, 7.8, True, NAVY)


def add_footer(slide, num: int, width, height):
    textbox(slide, Cm(0.7), height - Cm(0.62), width - Cm(1.4), Cm(0.26),
            f"해군교전분석모델 Build-II 결과분석 제안서    {num:02d}", 5.8, False, GRAY, PP_ALIGN.RIGHT)


def bullet_panel(slide, bullets: Sequence[str], width, height, max_items=4):
    y = height - Cm(4.05)
    textbox(slide, Cm(0.75), y, Cm(3.0), Cm(0.35), "핵심 포인트", 8, True, NAVY)
    sh = rect(slide, Cm(0.7), y + Cm(0.44), width - Cm(1.4), Cm(2.55), PALE, LINE)
    tf = sh.text_frame
    tf.clear()
    tf.word_wrap = True
    tf.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
    tf.margin_left = Cm(0.28)
    tf.margin_right = Cm(0.28)
    tf.margin_top = Cm(0.12)
    for i, b in enumerate(bullets[:max_items]):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = f"• {b}"
        p.font.name = "Malgun Gothic"
        p.font.size = Pt(7.2)
        p.font.color.rgb = DARK
        p.alignment = PP_ALIGN.LEFT
        p.space_after = Pt(1)


def process(slide, labels: Sequence[str], x, y, w, h, color=BLUE):
    gap = Cm(0.35)
    bw = (w - gap * (len(labels) - 1)) / len(labels)
    for i, lab in enumerate(labels):
        xx = x + i * (bw + gap)
        box(slide, xx, y, bw, h, lab, fill=WHITE, line=color, size=7.2, bold=True, color=NAVY)
        if i < len(labels) - 1:
            arrow(slide, xx + bw + Cm(0.07), y + h / 2 - Cm(0.18), Cm(0.22), Cm(0.36), color)


def matrix_table(slide, headers: Sequence[str], rows: Sequence[Sequence[str]], x, y, widths: Sequence, row_h=Cm(0.9)):
    for i, h in enumerate(headers):
        box(slide, x + sum(widths[:i]), y, widths[i], row_h, h, fill=NAVY, line=NAVY, size=6.9, bold=True, color=WHITE)
    for r, row in enumerate(rows):
        for c, val in enumerate(row):
            cell = box(slide, x + sum(widths[:c]), y + row_h * (r + 1), widths[c], row_h, val,
                       fill=PALE if r % 2 == 0 else WHITE, line=LINE, size=6.5, color=DARK)
            if c > 0:
                cell.text_frame.paragraphs[0].alignment = PP_ALIGN.LEFT
                cell.text_frame.margin_left = Cm(0.16)


def add_rfp_reference(slide, refs: dict[str, List[str]], kind: str, x, y, w, h):
    return


def example_strip(slide, title: str, labels: Sequence[str], x, y, w, h):
    rect(slide, x, y, w, h, WHITE, LINE)
    textbox(slide, x + Cm(0.2), y + Cm(0.15), w - Cm(0.4), Cm(0.32), title, 6.8, True, NAVY, PP_ALIGN.LEFT)
    gap = Cm(0.22)
    box_w = (w - Cm(0.45) - gap * (len(labels) - 1)) / len(labels)
    for i, label in enumerate(labels):
        xx = x + Cm(0.22) + i * (box_w + gap)
        fill = PALE if i % 2 == 0 else WHITE
        box(slide, xx, y + Cm(0.68), box_w, h - Cm(0.9), label, fill=fill, line=BLUE, size=6.1, bold=True, color=NAVY)


def analysis_mockup(slide, title: str, metrics: Sequence[Tuple[str, str]], x, y, w, h):
    rect(slide, x, y, w, h, WHITE, LINE)
    textbox(slide, x + Cm(0.25), y + Cm(0.18), w - Cm(0.5), Cm(0.35), title, 7.0, True, NAVY, PP_ALIGN.LEFT)
    map_x, map_y = x + Cm(0.35), y + Cm(0.75)
    map_w, map_h = w * 0.54, h - Cm(1.15)
    rect(slide, map_x, map_y, map_w, map_h, PALE, LINE)
    for gx in [0.25, 0.5, 0.75]:
        line(slide, map_x + map_w * gx, map_y + Cm(0.18), map_x + map_w * gx, map_y + map_h - Cm(0.18), LINE, 0.8)
    for gy in [0.33, 0.66]:
        line(slide, map_x + Cm(0.18), map_y + map_h * gy, map_x + map_w - Cm(0.18), map_y + map_h * gy, LINE, 0.8)
    box(slide, map_x + Cm(0.45), map_y + map_h - Cm(0.8), Cm(1.4), Cm(0.45), "플랫폼", fill=WHITE, line=BLUE, size=5.6, bold=True)
    box(slide, map_x + map_w - Cm(2.0), map_y + Cm(0.55), Cm(1.4), Cm(0.45), "표적", fill=WHITE, line=RED, size=5.6, bold=True, color=RED)
    line(slide, map_x + Cm(1.85), map_y + map_h - Cm(0.6), map_x + map_w - Cm(2.0), map_y + Cm(0.78), BLUE, 1.6)
    for i, color in enumerate([BLUE, GREEN, RED]):
        dot = slide.shapes.add_shape(MSO_SHAPE.OVAL, int(map_x + Cm(3.0 + i * 1.25)), int(map_y + Cm(1.0 + i * 0.4)), int(Cm(0.18)), int(Cm(0.18)))
        dot.fill.solid(); dot.fill.fore_color.rgb = color; dot.line.color.rgb = color
    panel_x = x + w * 0.60
    panel_w = w * 0.36
    for i, (name, val) in enumerate(metrics[:3]):
        yy = y + Cm(0.85 + i * 1.0)
        box(slide, panel_x, yy, panel_w, Cm(0.72), f"{name}\n{val}", fill=PALE if i % 2 == 0 else WHITE, line=BLUE, size=6.4, bold=True)


def decoy_mockup(slide, title: str, x, y, w, h):
    if "어뢰" in title:
        inputs = ["어뢰 위치/속도", "표적 기동", "기만기 운용모드", "수온·전달손실", "음선경로"]
        before = ["표적 추적 지속", "재획득 가능성 높음", "접근 경로 유지"]
        after = ["기만기 반응 반영", "추적 이탈/재획득 지연", "환경조건별 효과 비교"]
        outputs = ["기만 성공확률", "어뢰 추적 영향", "환경별 민감도"]
    else:
        inputs = ["유도탄 위치/속도", "표적 RCS/기동", "기만기 운용시점", "기만기 조합", "전자전 영향"]
        before = ["표적 직접 추적", "명중확률 기준값 산출", "위협 경로 유지"]
        after = ["RCS 변화 반영", "추적 혼선/오인 가능성", "기만 전후 확률 비교"]
        outputs = ["기만 성공확률", "명중확률 변화", "조건별 민감도"]

    rect(slide, x, y, w, h, WHITE, LINE)
    textbox(slide, x + Cm(0.25), y + Cm(0.18), w - Cm(0.5), Cm(0.35), title, 7.0, True, NAVY, PP_ALIGN.LEFT)

    col_y = y + Cm(0.75)
    col_h = h - Cm(1.05)
    col_gap = Cm(0.28)
    col_w = (w - Cm(0.7) - col_gap * 3) / 4
    xs = [x + Cm(0.35) + i * (col_w + col_gap) for i in range(4)]
    headers = ["입력조건", "기만 전", "기만 후", "산출결과"]
    columns = [inputs, before, after, outputs]
    fills = [PALE, WHITE, PALE, WHITE]
    for ci, (header, items) in enumerate(zip(headers, columns)):
        rect(slide, xs[ci], col_y, col_w, col_h, fills[ci], LINE)
        box(slide, xs[ci] + Cm(0.12), col_y + Cm(0.12), col_w - Cm(0.24), Cm(0.45),
            header, fill=NAVY if ci in (0, 3) else BLUE, line=BLUE, size=5.8, bold=True, color=WHITE)
        for ri, item in enumerate(items[:5]):
            yy = col_y + Cm(0.78 + ri * 0.58)
            box(slide, xs[ci] + Cm(0.18), yy, col_w - Cm(0.36), Cm(0.42),
                item, fill=WHITE, line=LINE, size=5.1, color=DARK, align=PP_ALIGN.LEFT)
    for ci in range(3):
        arrow(slide, xs[ci] + col_w + Cm(0.05), col_y + col_h / 2 - Cm(0.16), Cm(0.22), Cm(0.32), BLUE)


def title_slide(slide, width, height):
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = WHITE
    rect(slide, Cm(0.85), Cm(5.5), width - Cm(1.7), Cm(0.08), BLUE)
    textbox(slide, Cm(1.0), Cm(6.05), width - Cm(2), Cm(2.0), "결과분석 및 전장 가시화\n체계 구축 제안서", 23, True, DARK)
    textbox(slide, Cm(1.05), Cm(9.0), width - Cm(2.1), Cm(0.9), "해군교전분석모델 Build-II | 결과분석 고도화 수행방안", 10.5, False, NAVY)
    process(slide, ["모의결과", "분석엔진", "3D Replay", "Dashboard", "Report"], Cm(1.05), Cm(12.2), width - Cm(2.1), Cm(1.25))
    textbox(slide, Cm(1.05), height - Cm(1.6), width - Cm(2.1), Cm(0.4), "전투 결과를 의사결정 자료로 전환하는 결과분석 플랫폼 구축", 8.5, False, GRAY)


def section_slide(slide, width, height):
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = WHITE
    textbox(slide, Cm(1.0), Cm(6.5), Cm(2.8), Cm(1.6), "Ⅱ", 34, True, NAVY, PP_ALIGN.CENTER)
    textbox(slide, Cm(4.2), Cm(6.85), width - Cm(5.2), Cm(0.9), "전략 및 방법론", 18, True, DARK)
    rect(slide, Cm(4.25), Cm(7.85), width - Cm(5.4), Cm(0.05), BLUE)
    textbox(slide, Cm(4.25), Cm(8.25), width - Cm(5.4), Cm(1.1),
            "사업 이해 / 결과분석 수행방안 / 전장 가시화 / 시험평가 대응", 10, False, GRAY)


def architecture(slide, width, y_start=Cm(4.35)):
    layers = [
        ("전시 계층", ["3D 전장 가시화", "Multi View", "KPI Dashboard", "자동 보고서"]),
        ("분석 계층", ["탐지확률", "명중확률", "기만성공확률", "전자전 효과"]),
        ("데이터 계층", ["이벤트 로그", "플랫폼 상태", "무기체계 상태", "수중환경"]),
        ("연계/저장", ["분석 API", "DB 저장", "파일 저장소", "추적 이력"]),
    ]
    y = y_start
    for idx, (name, items) in enumerate(layers):
        box(slide, Cm(0.9), y, Cm(2.8), Cm(1.0), name, fill=NAVY if idx == 0 else BLUE, line=BLUE, size=7.2, bold=True, color=WHITE)
        item_w = (width - Cm(4.7)) / len(items)
        for i, item in enumerate(items):
            box(slide, Cm(4.0) + i * item_w, y, item_w - Cm(0.15), Cm(1.0), item, fill=PALE if idx % 2 == 0 else WHITE, line=LINE, size=6.4, bold=True)
        y += Cm(1.45)


def data_flow(slide, width):
    process(slide, ["시뮬레이션 실행", "이벤트 로그 수집", "분석 엔진 처리", "DB 저장", "3D Replay / Dashboard / Report"],
            Cm(0.9), Cm(4.5), width - Cm(1.8), Cm(1.25))
    rows = [
        ["수집", "시나리오, 플랫폼 상태, 무기 이벤트, 환경자료"],
        ["정제", "시간축 정렬, 객체 식별자 매핑, 누락/오류 검증"],
        ["분석", "확률모델 계산, 기만/전자전 효과 산출, KPI 집계"],
        ["활용", "Replay, 차트, 보고서, 시험평가 증적"],
    ]
    matrix_table(slide, ["단계", "처리 내용"], rows, Cm(1.0), Cm(7.0), [Cm(3.0), width - Cm(5.0)], Cm(0.9))


def db_schema(slide, width):
    tables = {
        "SIM_RUN": ["run_id", "scenario_id", "start_time", "version"],
        "PLATFORM_STATE": ["run_id", "object_id", "time", "position"],
        "WEAPON_EVENT": ["run_id", "weapon_id", "event_type", "target_id"],
        "ANALYSIS_RESULT": ["run_id", "metric_type", "time", "value"],
    }
    x_positions = [Cm(0.9), Cm(9.1), Cm(0.9), Cm(9.1)]
    y_positions = [Cm(4.3), Cm(4.3), Cm(8.4), Cm(8.4)]
    for idx, (name, fields) in enumerate(tables.items()):
        x, y = x_positions[idx], y_positions[idx]
        box(slide, x, y, Cm(6.6), Cm(0.65), name, fill=NAVY, line=NAVY, size=7, bold=True, color=WHITE)
        box(slide, x, y + Cm(0.72), Cm(6.6), Cm(2.3), "\n".join(fields), fill=PALE, line=LINE, size=6.5, color=DARK, align=PP_ALIGN.LEFT)
    line(slide, Cm(7.5), Cm(5.8), Cm(9.1), Cm(5.8))
    line(slide, Cm(7.5), Cm(9.9), Cm(9.1), Cm(9.9))
    line(slide, Cm(4.2), Cm(7.3), Cm(4.2), Cm(8.4))


def replay_diagram(slide, width):
    process(slide, ["로그 로딩", "타임라인 생성", "이벤트 상태 복원", "3D 장면 재생", "구간 북마크"],
            Cm(0.9), Cm(4.35), width - Cm(1.8), Cm(1.1))
    x, y, w, h = Cm(1.0), Cm(6.6), width - Cm(2.0), Cm(4.8)
    rect(slide, x, y, w, h, WHITE, LINE)
    rect(slide, x + Cm(0.35), y + Cm(0.35), w - Cm(0.7), Cm(2.55), PALE, LINE)
    box(slide, x + Cm(0.7), y + Cm(1.1), Cm(2.0), Cm(0.55), "함정", fill=WHITE, line=BLUE, size=6.5, bold=True)
    box(slide, x + w - Cm(3.0), y + Cm(1.65), Cm(2.0), Cm(0.55), "표적", fill=WHITE, line=BLUE, size=6.5, bold=True)
    line(slide, x + Cm(2.7), y + Cm(1.4), x + w - Cm(3.0), y + Cm(1.92), BLUE, 2)
    rect(slide, x + Cm(0.7), y + Cm(3.45), w - Cm(1.4), Cm(0.12), BLUE)
    for i, lab in enumerate(["00:00", "탐지", "발사", "기만", "명중/실패"]):
        xx = x + Cm(0.75) + i * (w - Cm(1.6)) / 4
        box(slide, xx - Cm(0.5), y + Cm(3.8), Cm(1.0), Cm(0.4), lab, fill=WHITE, line=LINE, size=5.4)


def multiview(slide, width):
    x, y, w, h = Cm(0.9), Cm(4.15), width - Cm(1.8), Cm(8.5)
    cells = [
        ("전장뷰\n3D 공간/궤적/이벤트", x, y, w * 0.58, h * 0.58, PALE),
        ("플랫폼뷰\n함정·항공기·잠수함 상태", x + w * 0.60, y, w * 0.40, h * 0.28, WHITE),
        ("무장뷰\n유도탄·어뢰·기만기", x + w * 0.60, y + h * 0.32, w * 0.40, h * 0.26, WHITE),
        ("이벤트 로그뷰\n시간순 이벤트 및 필터", x, y + h * 0.62, w, h * 0.38, WHITE),
    ]
    for text, cx, cy, cw, ch, fill in cells:
        box(slide, cx, cy, cw - Cm(0.1), ch - Cm(0.1), text, fill=fill, line=BLUE, size=7.2, bold=True)


def underwater(slide, width):
    x, y, w, h = Cm(0.9), Cm(4.2), width - Cm(1.8), Cm(7.3)
    rect(slide, x, y, w, h, WHITE, LINE)
    for i, color in enumerate([RGBColor(0xD9, 0xF0, 0xFF), RGBColor(0xB7, 0xDD, 0xEE), RGBColor(0x8B, 0xC0, 0xD8), RGBColor(0x5E, 0x9F, 0xBD)]):
        rect(slide, x + Cm(0.4), y + Cm(0.7 + i * 1.25), w - Cm(0.8), Cm(1.1), color, color)
    line(slide, x + Cm(1.1), y + Cm(5.8), x + w - Cm(1.0), y + Cm(1.2), BLUE, 2.3)
    line(slide, x + Cm(1.3), y + Cm(5.2), x + w - Cm(1.4), y + Cm(2.0), GREEN, 1.8)
    for i, lab in enumerate(["수온", "수중분포", "음파전달손실", "음선경로"]):
        box(slide, x + Cm(0.6 + i * 3.8), y + h - Cm(0.75), Cm(3.1), Cm(0.45), lab, fill=WHITE, line=BLUE, size=6.1, bold=True)


def probability_model(slide, width, title: str, outputs: Sequence[str]):
    inputs = ["거리", "속도", "RCS", "재머효과", "수중환경"]
    x, y = Cm(0.95), Cm(4.15)
    for i, inp in enumerate(inputs):
        box(slide, x, y + i * Cm(0.68), Cm(3.0), Cm(0.48), inp, fill=PALE, line=LINE, size=6.3, bold=True)
    model_x, model_y, model_w, model_h = Cm(5.25), Cm(5.15), Cm(4.05), Cm(1.45)
    box(slide, model_x, model_y, model_w, model_h, title, fill=NAVY, line=NAVY, size=8, bold=True, color=WHITE)
    model_center_y = model_y + model_h / 2
    for i in range(len(inputs)):
        in_center_y = y + i * Cm(0.68) + Cm(0.24)
        line(slide, x + Cm(3.0), in_center_y, model_x, model_center_y, BLUE, 1)
    for i, out in enumerate(outputs):
        out_y = Cm(4.45 + i * 1.18)
        out_h = Cm(0.78)
        out_x = Cm(10.75)
        out_center_y = out_y + out_h / 2
        box(slide, out_x, out_y, Cm(4.9), out_h, out, fill=WHITE, line=BLUE, size=7.1, bold=True)
        line(slide, model_x + model_w, model_center_y, Cm(9.65), out_center_y, BLUE, 1.1)
        arrow(slide, Cm(9.68), out_center_y - Cm(0.17), Cm(0.55), Cm(0.34), BLUE)
        line(slide, Cm(10.23), out_center_y, out_x, out_center_y, BLUE, 1.1)


def ew_before_after(slide, width):
    box(slide, Cm(1.0), Cm(4.4), Cm(6.9), Cm(4.2), "전자전 OFF\n\n탐지/추적 영향 미반영\n명중률 예시: 74%", fill=PALE, line=LINE, size=11, bold=True)
    box(slide, Cm(9.3), Cm(4.4), Cm(6.9), Cm(4.2), "전자전 ON\n\n재머·GPS 교란 영향 반영\n명중률 예시: 39%", fill=WHITE, line=BLUE, size=11, bold=True)
    arrow(slide, Cm(8.05), Cm(6.1), Cm(1.0), Cm(0.55), BLUE)
    matrix_table(slide, ["효과 항목", "분석 반영 내용"], [
        ["재머", "탐지/추적 성능 저하 효과를 이벤트 시간축에 반영"],
        ["GPS 교란", "항법 오차 및 유도 영향도를 교전결과와 연계"],
        ["전자전 효과", "교전 전후 확률 변화와 영향 대상을 비교"],
    ], Cm(1.0), Cm(9.4), [Cm(4.0), width - Cm(6.0)], Cm(0.78))


def kpi_dashboard(slide, width):
    metrics = [("탐지확률", "82%"), ("명중확률", "64%"), ("기만성공", "41%"), ("전자전효과", "중")]
    for i, (name, val) in enumerate(metrics):
        box(slide, Cm(1.0 + i * 3.95), Cm(4.2), Cm(3.4), Cm(1.5), f"{name}\n{val}", fill=PALE if i % 2 == 0 else WHITE, line=BLUE, size=10, bold=True)
    probability_chart(slide, Cm(1.0), Cm(6.5), width - Cm(2.0), Cm(5.2), "KPI 추이 및 이벤트 연계")


def probability_chart(slide, x, y, w, h, title: str):
    rect(slide, x, y, w, h, WHITE, LINE)
    textbox(slide, x + Cm(0.35), y + Cm(0.25), w - Cm(0.7), Cm(0.35), title, 7.5, True, NAVY)
    left, top, cw, ch = x + Cm(0.6), y + Cm(1.0), w - Cm(1.1), h - Cm(1.6)
    for i in range(4):
        rect(slide, left, top + ch * i / 3, cw, Cm(0.01), LINE, LINE)
    last = None
    for px, py in [(0, .72), (.16, .61), (.34, .48), (.55, .56), (.72, .34), (1, .42)]:
        sx, sy = left + cw * px, top + ch * py
        dot = slide.shapes.add_shape(MSO_SHAPE.OVAL, int(sx - Cm(0.04)), int(sy - Cm(0.04)), int(Cm(0.08)), int(Cm(0.08)))
        dot.fill.solid(); dot.fill.fore_color.rgb = BLUE; dot.line.color.rgb = BLUE
        if last:
            line(slide, last[0], last[1], sx, sy, BLUE, 1.7)
        last = (sx, sy)


SLIDES = [
    ("표지", "", "cover"),
    ("사업 이해", "Build-II 결과분석은 모의결과의 해석성, 재현성, 검증성을 강화하는 체계 구축이 핵심입니다.", "understanding"),
    ("결과분석 업무 범위", "결과분석 담당 범위를 데이터, 모델, 가시화, 리포트 산출까지 명확히 정의합니다.", "scope"),
    ("결과분석 체계 개념도", "전투 결과를 의사결정 자료로 전환하는 결과분석 플랫폼 구조를 제안합니다.", "concept"),
    ("전체 시스템 아키텍처", "데이터 계층, 분석 계층, 전시 계층을 모듈화하여 확장 가능한 구조를 구성합니다.", "arch"),
    ("데이터 수집 및 처리 흐름", "모의 실행 결과를 표준 처리 흐름으로 정제하여 Replay, Dashboard, Report로 활용합니다.", "flow"),
    ("결과분석 DB 구조 개념", "시뮬레이션 실행, 플랫폼 상태, 무기 이벤트, 분석 결과를 추적 가능한 DB 구조로 관리합니다.", "db"),
    ("3D 전장 가시화 개선 방안", "전장상황, 플랫폼, 무기체계, 결과분석 데이터를 동일 장면에서 통합 전시합니다.", "3d"),
    ("자동 3D Replay 구성", "모의결과 로그를 기반으로 시간축 Replay와 이벤트 구간 분석을 제공합니다.", "replay"),
    ("Multi View 화면 구성", "전장뷰, 플랫폼뷰, 무장뷰, 이벤트 로그뷰를 4분할 화면으로 연계합니다.", "multi"),
    ("수중환경 3D 가시화", "수온, 분포, 전달손실, 음선경로를 수중환경 레이어로 시각화합니다.", "underwater"),
    ("음파전달손실/음선경로 가시화", "음향환경 변화가 어뢰 및 기만 분석에 미치는 영향을 직관적으로 표현합니다.", "sound"),
    ("계산모델 가시화", "입력값, 계산 단계, 중간 산출값, 출력값을 추적 가능한 형태로 제시합니다.", "model"),
    ("탐지확률 실시간 차트", "거리, 속도, RCS, 전자전 효과 등 입력조건을 기반으로 탐지확률을 시계열로 표시합니다.", "detect"),
    ("명중확률 실시간 차트", "교전 조건과 무기체계 상태를 반영하여 명중확률 변화를 실시간 분석합니다.", "hit"),
    ("유도탄-기만기 기만 성공확률 개선", "유도탄 추적 상태와 기만기 운용 조건을 반영하여 기만 성공확률을 산출합니다.", "missile"),
    ("어뢰-기만기 기만 성공확률 개선", "수중환경과 음향 조건을 반영하여 어뢰 기만 효과를 분석합니다.", "torpedo"),
    ("전자전 장비 효과 산출", "재머, GPS 교란 등 전자전 효과를 교전 결과와 연계하여 비교 분석합니다.", "ew"),
    ("KPI Dashboard", "핵심 분석 지표를 요약하여 지휘관/분석관의 신속한 판단을 지원합니다.", "kpi"),
    ("자동 보고서 생성", "Replay, 차트, 이벤트 요약을 자동 보고서로 연계하여 분석 산출물 품질을 높입니다.", "report"),
    ("기술 스택", "검증 가능한 표준 기술 조합으로 구현 가능성과 운영성을 확보합니다.", "tech"),
    ("개발 일정/WBS", "분석, 설계, 구현, 통합, 시험평가를 단계적으로 추진합니다.", "wbs"),
    ("위험요소 및 대응방안", "성능, 데이터 정합성, 가시화 복잡도, 시험평가 리스크를 사전에 관리합니다.", "risk"),
    ("기대효과", "전술분석 고도화, 무기체계 효과도 검증, 교전결과 이해도 향상을 기대합니다.", "effect"),
    ("결론", "Build-II 요구사항에 부합하는 결과분석 및 전장 가시화 플랫폼을 단계적으로 구축합니다.", "conclusion"),
]


def slide_body(slide, kind: str, width, height, rfp_refs: dict[str, List[str]]):
    if kind == "understanding":
        process(slide, ["RFP 요구", "결과분석 범위", "가시화/Replay", "확률·효과 분석", "시험평가"], Cm(0.9), Cm(4.3), width - Cm(1.8), Cm(1.0))
        matrix_table(slide, ["이해 관점", "제안 반영 방향"], [
            ["결과분석", "모의결과를 단순 결과표가 아닌 분석 가능한 데이터 자산으로 전환"],
            ["전장 가시화", "전장환경, 플랫폼, 무기체계의 상호작용을 3D 및 멀티뷰로 표현"],
            ["확률/효과 분석", "탐지·명중·기만·전자전 효과를 시간축 기반 차트로 분석"],
            ["시험평가", "요구사항별 산출물과 시험 기준을 연결하여 검증 가능성 확보"],
        ], Cm(1.0), Cm(6.5), [Cm(4.0), width - Cm(6.0)], Cm(0.85))
        add_rfp_reference(slide, rfp_refs, "understanding", Cm(0.9), Cm(10.15), width - Cm(1.8), Cm(2.25))
        example_strip(slide, "제안 대응 구조", ["요구사항 분석", "기능 설계", "구현 검증", "운영 활용"], Cm(0.9), Cm(10.2), width - Cm(1.8), Cm(2.0))
    elif kind == "scope":
        matrix_table(slide, ["구분", "담당 범위", "주요 산출물"], [
            ["데이터", "로그 수집, 정제, 객체/시간 매핑", "분석 DB, 정합성 점검 결과"],
            ["분석모델", "탐지·명중·기만·전자전 산출", "확률/효과 분석값"],
            ["가시화", "3D Replay, Multi View, Chart", "분석 화면, 캡처, 이벤트 뷰"],
            ["보고", "KPI 및 자동 보고서 구성", "보고서 템플릿, 시험 증적"],
        ], Cm(0.9), Cm(4.3), [Cm(3.1), Cm(6.5), width - Cm(11.4)], Cm(0.95))
        example_strip(slide, "결과분석 업무 산출 흐름", ["로그 정제", "분석값 산출", "3D/차트 전시", "보고서 증적화"], Cm(0.9), Cm(9.65), width - Cm(1.8), Cm(2.0))
    elif kind == "concept":
        process(slide, ["모의결과", "분석 DB", "분석 엔진", "3D/Chart", "보고서"], Cm(0.9), Cm(4.4), width - Cm(1.8), Cm(1.25))
        architecture(slide, width, y_start=Cm(6.35))
    elif kind == "arch":
        architecture(slide, width)
        add_rfp_reference(slide, rfp_refs, "arch", Cm(0.9), Cm(10.45), width - Cm(1.8), Cm(2.1))
        example_strip(slide, "아키텍처 연계 예시", ["모의 로그", "분석 API", "결과분석 DB", "전시/보고"], Cm(0.9), Cm(10.55), width - Cm(1.8), Cm(2.0))
    elif kind == "flow":
        data_flow(slide, width)
        add_rfp_reference(slide, rfp_refs, "flow", Cm(1.0), Cm(12.15), width - Cm(2.0), Cm(2.05))
        example_strip(slide, "처리 결과 활용 예시", ["시계열 정렬", "이벤트 복원", "확률 집계", "분석 리포트"], Cm(1.0), Cm(12.15), width - Cm(2.0), Cm(1.85))
    elif kind == "db":
        db_schema(slide, width)
        example_strip(slide, "DB 저장 후 활용 예시", ["Replay 재현", "KPI 집계", "조건별 비교", "보고서 추적"], Cm(1.0), Cm(12.15), width - Cm(2.0), Cm(1.75))
    elif kind == "3d":
        process(slide, ["객체 표준화", "좌표/시간 동기화", "분석값 오버레이", "3D 장면 저장"], Cm(0.9), Cm(4.3), width - Cm(1.8), Cm(1.1))
        matrix_table(slide, ["구현 항목", "내용"], [
            ["객체 표현", "함정, 항공기, 잠수함, 유도탄, 어뢰, 기만기 상태를 3D 객체로 표현"],
            ["분석 오버레이", "탐지·명중·기만·전자전 결과를 색상, 라벨, 궤적으로 표시"],
            ["활용", "Replay 캡처, 분석 장면 저장, 보고서 연계"],
        ], Cm(1.0), Cm(6.6), [Cm(3.4), width - Cm(5.4)], Cm(0.9))
        add_rfp_reference(slide, rfp_refs, "3d", Cm(0.9), Cm(10.75), width - Cm(1.8), Cm(2.15))
        analysis_mockup(slide, "3D 가시화 화면 구성 예시", [("객체", "플랫폼"), ("궤적", "시간축"), ("분석값", "오버레이")], Cm(0.9), Cm(10.7), width - Cm(1.8), Cm(4.2))
    elif kind == "replay":
        replay_diagram(slide, width)
    elif kind == "multi":
        multiview(slide, width)
    elif kind in ("underwater", "sound"):
        underwater(slide, width)
        if kind == "sound":
            matrix_table(slide, ["분석 항목", "가시화 방식"], [
                ["음파전달손실", "거리/심도별 손실 값을 색상 레이어로 표시"],
                ["음선경로", "음파 굴절 경로를 곡선으로 표시"],
                ["분석 연계", "어뢰 탐지 및 기만 성공확률 산출 입력으로 활용"],
            ], Cm(1.0), Cm(11.9), [Cm(4.0), width - Cm(6.0)], Cm(0.75))
        ref_y = Cm(15.35) if kind == "sound" else Cm(13.0)
        add_rfp_reference(slide, rfp_refs, kind, Cm(0.9), ref_y, width - Cm(1.8), Cm(2.05))
        if kind == "underwater":
            example_strip(slide, "수중환경 분석 연계", ["수온분포", "전달손실", "음선경로", "기만확률 입력"], Cm(0.9), Cm(12.8), width - Cm(1.8), Cm(2.0))
    elif kind == "model":
        process(slide, ["입력값", "모델 구조", "중간 산출", "출력값", "근거 리포트"], Cm(0.9), Cm(4.4), width - Cm(1.8), Cm(1.15))
        matrix_table(slide, ["표시 항목", "제안 내용"], [
            ["입력 파라미터", "거리, 속도, 환경, 센서, 무기체계 상태를 조회"],
            ["계산 흐름", "모델별 산출 단계와 중간값을 추적"],
            ["결과 근거", "산출식 버전과 입력값을 결과 리포트에 포함"],
        ], Cm(1.0), Cm(6.7), [Cm(3.6), width - Cm(5.6)], Cm(0.85))
        add_rfp_reference(slide, rfp_refs, "model", Cm(1.0), Cm(10.3), width - Cm(2.0), Cm(2.05))
        example_strip(slide, "계산모델 추적 예시", ["입력값", "계산식/버전", "중간값", "산출근거"], Cm(1.0), Cm(10.3), width - Cm(2.0), Cm(2.0))
    elif kind == "detect":
        probability_model(slide, width, "탐지확률 계산 모델", ["탐지확률", "탐지 이벤트", "임계값 비교"])
        probability_chart(slide, Cm(1.0), Cm(10.2), width - Cm(2), Cm(3.0), "탐지확률 시계열")
        analysis_mockup(slide, "탐지확률 분석 화면 예시", [("탐지확률", "82%"), ("임계값", "70%"), ("상태", "탐지")], Cm(1.0), Cm(13.7), width - Cm(2.0), Cm(5.0))
    elif kind == "hit":
        probability_model(slide, width, "명중확률 계산 모델", ["명중확률", "교전 이벤트", "조건별 비교"])
        probability_chart(slide, Cm(1.0), Cm(10.2), width - Cm(2), Cm(3.0), "명중확률 시계열")
        analysis_mockup(slide, "명중확률 분석 화면 예시", [("명중확률", "64%"), ("교전단계", "중간유도"), ("판정", "비교")], Cm(1.0), Cm(13.7), width - Cm(2.0), Cm(5.0))
    elif kind == "missile":
        probability_model(slide, width, "유도탄-기만 효과 모델", ["기만 성공확률", "기만 전/후 비교", "민감도 분석"])
        decoy_mockup(slide, "유도탄 기만효과 비교 예시", Cm(1.0), Cm(9.6), width - Cm(2.0), Cm(5.2))
    elif kind == "torpedo":
        probability_model(slide, width, "어뢰-기만 효과 모델", ["기만 성공확률", "환경별 비교", "어뢰 추적 영향"])
        decoy_mockup(slide, "어뢰 기만효과 비교 예시", Cm(1.0), Cm(9.6), width - Cm(2.0), Cm(5.2))
    elif kind == "ew":
        ew_before_after(slide, width)
        add_rfp_reference(slide, rfp_refs, "ew", Cm(1.0), Cm(12.85), width - Cm(2.0), Cm(2.05))
        example_strip(slide, "전자전 효과 비교 관점", ["재머 영향", "GPS 교란", "확률 변화", "효과 판정"], Cm(1.0), Cm(12.85), width - Cm(2.0), Cm(1.9))
    elif kind == "kpi":
        kpi_dashboard(slide, width)
    elif kind == "report":
        process(slide, ["이벤트 요약", "차트 캡처", "Replay 장면", "KPI 표", "보고서 생성"], Cm(0.9), Cm(4.4), width - Cm(1.8), Cm(1.1))
        matrix_table(slide, ["보고서 구성", "내용"], [
            ["분석 개요", "시나리오, 모의조건, 분석 대상"],
            ["핵심 결과", "탐지·명중·기만·전자전 지표"],
            ["근거 자료", "이벤트 로그, Replay 장면, 차트 캡처"],
            ["시험 증적", "요구사항별 확인 결과와 판정 근거"],
        ], Cm(1.0), Cm(6.6), [Cm(3.5), width - Cm(5.5)], Cm(0.78))
        example_strip(slide, "자동 보고서 구성 예시", ["요약 표지", "확률 차트", "Replay 캡처", "판정 근거"], Cm(1.0), Cm(10.6), width - Cm(2.0), Cm(1.85))
    elif kind == "tech":
        matrix_table(slide, ["영역", "적용 기술 방향"], [
            ["3D/가시화", "3D 엔진 기반 전장 객체, 궤적, 이벤트 오버레이"],
            ["데이터", "결과분석 DB, 로그 정제/매핑 처리"],
            ["분석", "확률 계산 모델, 기만/전자전 효과 산출 모듈"],
            ["UI", "멀티뷰, 차트, 대시보드, 보고서 생성 화면"],
        ], Cm(1.0), Cm(4.4), [Cm(3.6), width - Cm(5.6)], Cm(0.9))
        example_strip(slide, "구현 구성 예시", ["3D Viewer", "Analysis API", "Result DB", "Dashboard"], Cm(1.0), Cm(9.1), width - Cm(2.0), Cm(2.0))
    elif kind == "wbs":
        matrix_table(slide, ["단계", "주요 수행 내용", "산출물"], [
            ["분석", "요구사항 정제, 데이터 항목 정의", "요구사항 매트릭스"],
            ["설계", "DB, 분석모델, UI/가시화 구조 설계", "설계서/화면정의서"],
            ["구현", "Replay, Multi View, 확률/효과 분석 구현", "구현 모듈"],
            ["통합", "모의 로그 연계, DB 저장, 화면 통합", "통합 결과"],
            ["시험", "시나리오 기반 기능/정합성 시험", "시험결과서"],
        ], Cm(0.9), Cm(4.3), [Cm(2.6), Cm(7.2), width - Cm(11.6)], Cm(0.78))
        example_strip(slide, "단계별 중점 관리", ["요구 추적", "설계 검토", "통합 점검", "시험 증적"], Cm(0.9), Cm(9.4), width - Cm(1.8), Cm(1.9))
    elif kind == "risk":
        matrix_table(slide, ["위험요소", "영향", "대응방안"], [
            ["로그 정합성 부족", "Replay 재현성 저하", "정합성 점검 규칙과 예외 처리 기준 수립"],
            ["3D 성능 저하", "대용량 시나리오 표시 지연", "LOD, 필터링, 구간 로딩 적용"],
            ["모델 산출 근거 부족", "시험평가 설명력 저하", "입력값-산출값 추적 정보 저장"],
            ["화면 복잡도 증가", "분석관 사용성 저하", "멀티뷰 템플릿과 역할별 화면 구성"],
        ], Cm(0.9), Cm(4.3), [Cm(3.6), Cm(4.4), width - Cm(9.8)], Cm(0.85))
        example_strip(slide, "리스크 관리 관점", ["데이터 정합성", "성능", "설명가능성", "사용성"], Cm(0.9), Cm(9.0), width - Cm(1.8), Cm(1.9))
    elif kind == "effect":
        boxes = ["전술분석 고도화", "무기체계 효과도 검증", "교전결과 직관적 이해", "의사결정 지원", "교육/훈련 활용"]
        process(slide, boxes, Cm(0.9), Cm(4.7), width - Cm(1.8), Cm(1.15))
        matrix_table(slide, ["기대효과", "설명"], [
            ["해석성", "모의결과를 3D, 차트, 리포트로 다각도 해석"],
            ["재현성", "동일 로그 기반 Replay로 분석 결과 반복 확인"],
            ["검증성", "입력값, 모델, 산출값의 추적 구조 확보"],
            ["운영성", "분석관 업무 흐름에 맞춘 대시보드와 보고서 제공"],
        ], Cm(1.0), Cm(7.2), [Cm(3.4), width - Cm(5.4)], Cm(0.8))
        example_strip(slide, "활용 장면", ["전술 검토", "효과도 비교", "시험평가 설명", "교육자료 전환"], Cm(1.0), Cm(11.2), width - Cm(2.0), Cm(1.85))
    elif kind == "conclusion":
        matrix_table(slide, ["구분", "최종 제안 메시지"], [
            ["결과분석 고도화", "모의결과의 해석성, 재현성, 검증성을 강화"],
            ["3D 전장 가시화", "자동 Replay와 Multi View로 분석 효율 향상"],
            ["확률/효과 분석", "탐지·명중·기만·전자전 효과를 통합 분석"],
            ["확장 기반", "향후 분석모델과 화면 기능 확장이 가능한 플랫폼 기반 확보"],
        ], Cm(1.0), Cm(4.6), [Cm(4.0), width - Cm(6.0)], Cm(0.9))
        example_strip(slide, "제안 방향 요약", ["요구 충족", "구현 가능성", "검증 중심", "확장 기반"], Cm(1.0), Cm(9.45), width - Cm(2.0), Cm(2.0))


def bullets_for(kind: str) -> List[str]:
    return {
        "understanding": ["결과분석을 의사결정 지원 자료로 전환", "전장상황과 확률 산출 결과를 연계", "시험평가 설명력을 확보"],
        "scope": ["데이터-모델-가시화-보고 범위를 명확화", "결과분석 담당 산출물 기준 수립", "검증 가능한 산출 흐름 구성"],
        "concept": ["모의결과 저장 후 분석/전시/보고로 연계", "전장 가시화와 수치 분석을 통합", "분석 리포트 생성까지 일관 처리"],
        "arch": ["계층별 모듈화로 확장성 확보", "분석모델과 UI를 분리하여 유지관리성 강화", "표준 인터페이스 기반 연계"],
        "flow": ["로그 수집부터 보고서까지 단일 흐름 구성", "시간축 정렬로 Replay 재현성 확보", "DB 저장 기반 추적성 확보"],
        "db": ["실행 단위 중심 데이터 구조", "상태/이벤트/결과 분리 저장", "분석 결과 재현을 위한 이력 관리"],
        "3d": ["3D 객체와 분석값을 동시 표현", "전장상황과 무기 이벤트를 시간축으로 연계", "장면 저장 및 보고서 연계 가능"],
        "replay": ["게임 Replay처럼 시간축 기반 재생", "이벤트 점프와 구간 반복 지원", "동일 로그 기반 재현성 확보"],
        "multi": ["분석관이 전체-상세-로그를 동시에 확인", "객체 선택 시 뷰 간 동기화", "플랫폼/무장/이벤트 분석 효율 향상"],
        "underwater": ["수중환경을 분석 레이어로 구성", "음향 정보와 어뢰/기만 분석 연계", "환경조건별 비교 분석 지원"],
        "sound": ["음파전달손실과 음선경로를 직관화", "수중 무기체계 분석 입력으로 활용", "환경 변화 영향 비교 가능"],
        "model": ["입력값과 산출값의 근거를 추적", "계산모델 구조를 분석관에게 공개", "결과 리포트 설명력 강화"],
        "detect": ["센서/환경 조건 기반 탐지확률 산출", "3D 객체와 차트 동기화", "임계값 및 이벤트 표시"],
        "hit": ["교전조건 기반 명중확률 산출", "발사-유도-결과 흐름과 연계", "조건별 비교 차트 제공"],
        "missile": ["유도탄 상태와 기만 이벤트 연계", "기만 성공확률 산출 절차 개선", "기만 전후 결과 비교"],
        "torpedo": ["수중환경과 어뢰 상태를 함께 반영", "기만기 효과를 환경조건별 비교", "어뢰 분석 설명력 강화"],
        "ew": ["전자전 OFF/ON 비교 구조 제공", "재머와 GPS 교란 효과 산출", "교전결과 변화와 연계"],
        "kpi": ["핵심 지표를 한 화면에서 확인", "이벤트와 지표 변화를 함께 해석", "보고서 연계 가능"],
        "report": ["Replay 장면과 차트를 자동 보고서화", "시험평가 증적 정리 시간 단축", "분석 결과 표준화"],
        "tech": ["구현 가능한 표준 기술 조합", "분석모델과 화면 기능 분리", "운영 확장성 확보"],
        "wbs": ["분석-설계-구현-통합-시험 단계화", "핵심 기능 우선 구현", "시험 기준과 산출물 연계"],
        "risk": ["기술/데이터/성능 리스크 사전 식별", "구현 단계별 대응방안 수립", "시험평가 리스크 완화"],
        "effect": ["전술분석 고도화", "교전결과 직관적 이해", "의사결정 및 교육훈련 활용"],
        "conclusion": ["Build-II 요구사항 중심의 구현 접근", "분석 효율과 검증성 동시 강화", "확장 가능한 결과분석 플랫폼 확보"],
    }.get(kind, [])


def build() -> None:
    task = read_task()
    rfp_refs = extract_rfp_references()
    prs = Presentation()
    w, h = slide_size()
    prs.slide_width = w
    prs.slide_height = h
    blank = prs.slide_layouts[6]

    for idx, (title, msg, kind) in enumerate(SLIDES, start=1):
        slide = prs.slides.add_slide(blank)
        if kind == "cover":
            title_slide(slide, w, h)
            continue
        if idx == 2:
            section_slide(slide, w, h)
            continue
        add_header(slide, idx, title, msg, w, h)
        slide_body(slide, kind, w, h, rfp_refs)
        bullet_panel(slide, bullets_for(kind), w, h)
        add_footer(slide, idx, w, h)
    prs.save(OUT)


if __name__ == "__main__":
    build()
