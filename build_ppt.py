# -*- coding: utf-8 -*-
from __future__ import annotations

import argparse
import csv
import math
import zipfile
import xml.etree.ElementTree as ET
from dataclasses import dataclass
from pathlib import Path
from typing import Iterable, List, Optional

import fitz
from PIL import Image, ImageDraw, ImageFilter, ImageOps, ImageStat
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR, MSO_AUTO_SIZE
from pptx.util import Cm, Pt


SCRIPT_DIR = Path(__file__).resolve().parent
ROOT = SCRIPT_DIR if (SCRIPT_DIR / "task.txt").exists() else SCRIPT_DIR.parent
OUTPUT = ROOT / "output"
ASSETS = OUTPUT / "assets"
RFP_PDF = ROOT / "해군교전2요청서.pdf"
TASK_TXT = ROOT / "task.txt"
FORM_DIR = ROOT / "form"
FORM_PPTX = FORM_DIR / "정성제안서.pptx"
FORM_WORK_PPTX = FORM_DIR / "3장_사업수행_계획 - 복사본.pptx"
FORM_PDF = FORM_DIR / "26AWAM정성제안서 1식).pdf"
PPTX_OUT = OUTPUT / "해군교전분석모델_BuildII_결과분석_제안서_초안.pptx"


NAVY = RGBColor(0x0B, 0x69, 0x93)
BLUE = RGBColor(0x11, 0xA0, 0xDD)
SKY = RGBColor(0x9D, 0xDD, 0xF9)
DARK = RGBColor(0x1F, 0x29, 0x33)
GRAY = RGBColor(0x66, 0x66, 0x66)
LIGHT = RGBColor(0xF4, 0xF7, 0xFA)
LINE = RGBColor(0xD8, 0xE2, 0xEA)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
PALE_BLUE = RGBColor(0xE9, 0xF5, 0xFB)
PALE_GRAY = RGBColor(0xF7, 0xF7, 0xF7)
WATERMARK = RGBColor(0xEA, 0xF2, 0xF7)


@dataclass
class Requirement:
    area: str
    item: str
    response: str
    verification: str
    priority: str = "상"


@dataclass
class SlidePlan:
    no: int
    title: str
    purpose: str
    message: str
    visual: str
    bullets: List[str]
    image_category: Optional[str] = None


@dataclass
class ImageCandidate:
    page: int
    category: str
    title: str
    path: Path
    score: float
    image_count: int
    drawing_count: int
    note: str


def ensure_dirs() -> None:
    OUTPUT.mkdir(exist_ok=True)
    ASSETS.mkdir(exist_ok=True)
    for stale in ASSETS.glob("rfp_p*.png"):
        stale.unlink()
    clean_template = OUTPUT / "_clean_정성제안서_template.pptx"
    if clean_template.exists():
        clean_template.unlink()


def clean_form_template() -> Optional[Path]:
    if not FORM_PPTX.exists():
        return None
    clean_path = OUTPUT / "_clean_정성제안서_template.pptx"
    ns_ct = {"ct": "http://schemas.openxmlformats.org/package/2006/content-types"}
    ns_pr = {
        "p": "http://schemas.openxmlformats.org/presentationml/2006/main",
        "r": "http://schemas.openxmlformats.org/officeDocument/2006/relationships",
    }
    ns_rel = {"rel": "http://schemas.openxmlformats.org/package/2006/relationships"}

    remove_terms = [
        "2025 태극 JOS 모델 유지보수 외주용역",
        "2025 태극JOS 모델 유지보수 외주용역",
        "2025",
        "태극",
        "태극 JOS",
        "태극JOS",
        "JOS",
        "유지보수",
        "외주용역",
    ]

    def sanitize_xml(data: bytes) -> bytes:
        if not data.lstrip().startswith(b"<"):
            return data
        text = data.decode("utf-8", errors="ignore")
        for term in remove_terms:
            text = text.replace(term, "")
        return text.encode("utf-8")

    with zipfile.ZipFile(FORM_PPTX, "r") as zin, zipfile.ZipFile(clean_path, "w", zipfile.ZIP_DEFLATED) as zout:
        for item in zin.infolist():
            name = item.filename
            if name.startswith("ppt/slides/") or name.startswith("ppt/notesSlides/"):
                continue
            data = zin.read(name)
            if name == "[Content_Types].xml":
                root = ET.fromstring(data)
                for override in list(root):
                    part = override.attrib.get("PartName", "")
                    if part.startswith("/ppt/slides/") or part.startswith("/ppt/notesSlides/"):
                        root.remove(override)
                data = ET.tostring(root, encoding="utf-8", xml_declaration=True)
            elif name == "ppt/presentation.xml":
                root = ET.fromstring(data)
                sld_id_lst = root.find("p:sldIdLst", ns_pr)
                if sld_id_lst is not None:
                    for child in list(sld_id_lst):
                        sld_id_lst.remove(child)
                data = ET.tostring(root, encoding="utf-8", xml_declaration=True)
            elif name == "ppt/_rels/presentation.xml.rels":
                root = ET.fromstring(data)
                for rel in list(root):
                    target = rel.attrib.get("Target", "")
                    typ = rel.attrib.get("Type", "")
                    if target.startswith("slides/") or typ.endswith("/slide"):
                        root.remove(rel)
                data = ET.tostring(root, encoding="utf-8", xml_declaration=True)
            else:
                data = sanitize_xml(data)
            zout.writestr(item, data)
    return clean_path


def read_task() -> str:
    raw = TASK_TXT.read_bytes()
    for enc in ("utf-8", "cp949", "euc-kr"):
        try:
            return raw.decode(enc)
        except UnicodeDecodeError:
            continue
    return raw.decode("utf-8", errors="replace")


def requirements() -> List[Requirement]:
    return [
        Requirement("결과분석 구조", "결과분석 데이터 3D 가시화 개선", "전장상황, 플랫폼, 무기체계, 분석 데이터를 동일 시간축과 좌표계로 연계 전시", "시나리오별 3D 재현, 객체 속성 조회, 화면 캡처 결과 확인"),
        Requirement("결과분석 구조", "모의결과 기반 자동 3D 재생", "시뮬레이션 로그와 이벤트를 시간축으로 정렬하여 재생, 일시정지, 배속, 구간 반복 제공", "로그-화면 동기화 점검, 이벤트 발생 시각 재현성 확인"),
        Requirement("결과분석 구조", "멀티뷰 기능", "전장환경, 함정, 항공기, 잠수함, 무기체계를 동시 뷰로 구성하고 상호작용 전시", "뷰 간 선택 객체 동기화, 필터/확대/추적 기능 시험"),
        Requirement("결과분석 구조", "수중환경 3D 가시화", "수온, 수중분포, 음파전달손실, 음선경로를 3D 레이어로 표현", "입력 환경자료와 렌더링 결과의 항목별 대조"),
        Requirement("결과분석 구조", "계산모델 가시화", "계산 흐름, 모델 구조, 입력값, 출력값을 추적 가능한 분석 뷰로 구성", "입출력 파라미터 추적, 산출 단계별 로그 대조"),
        Requirement("결과분석 구조", "실시간 분석 차트", "탐지확률, 명중확률, 기만성공확률, 전자전 효과를 시간축 차트로 제공", "시간 동기화, 산출식 버전, 임계값 표시 확인"),
        Requirement("결과분석 구성", "유도탄 기만 성공확률 개선", "기만기 운용 조건, 표적/유도탄 상태, 교전 이벤트를 반영하여 산출 절차 고도화", "대표 시나리오 회귀시험 및 민감도 분석"),
        Requirement("결과분석 구성", "어뢰 기만 성공확률 개선", "수중환경, 음향 조건, 기만기 효과를 반영한 어뢰 기만 확률 산출 개선", "환경조건별 산출값 비교 및 시험 기준 검증"),
        Requirement("결과분석 구성", "전자전 효과 산출 추가", "재머, GPS 교란 등 전자전 장비 효과를 탐지/교전 결과와 연계 분석", "전자전 이벤트-효과 산출-차트 표시 연계 확인"),
        Requirement("개발/검증", "시험평가 대응", "요구사항별 시험항목, 검증 기준, 추적 매트릭스를 수립", "요구사항-설계-구현-시험 추적성 검토"),
    ]


def render_pdf_assets() -> List[ImageCandidate]:
    doc = fitz.open(RFP_PDF)
    page_metrics = []
    thumbs = []
    for idx, page in enumerate(doc):
        pix = page.get_pixmap(matrix=fitz.Matrix(0.25, 0.25), alpha=False)
        img = Image.frombytes("RGB", [pix.width, pix.height], pix.samples)
        gray = ImageOps.grayscale(img)
        edge = gray.filter(ImageFilter.FIND_EDGES)
        score = ImageStat.Stat(edge).mean[0]
        page_metrics.append(
            {
                "page": idx + 1,
                "text_chars": len(page.get_text()),
                "images": len(page.get_images(full=True)),
                "drawings": len(page.get_drawings()),
                "edge_score": round(score, 3),
            }
        )
        t = img.resize((140, 198))
        d = ImageDraw.Draw(t)
        d.rectangle((0, 0, 42, 18), fill="white")
        d.text((3, 2), str(idx + 1), fill="black")
        thumbs.append(t)

    with (OUTPUT / "pdf_page_analysis.csv").open("w", encoding="utf-8-sig", newline="") as f:
        writer = csv.DictWriter(f, fieldnames=["page", "text_chars", "images", "drawings", "edge_score"])
        writer.writeheader()
        writer.writerows(page_metrics)

    cols = 5
    rows = math.ceil(len(thumbs) / cols)
    sheet = Image.new("RGB", (cols * 140, rows * 198), "white")
    for i, t in enumerate(thumbs):
        sheet.paste(t, ((i % cols) * 140, (i // cols) * 198))
    sheet.save(ASSETS / "rfp_contact_sheet.jpg", quality=92)

    candidates: List[ImageCandidate] = []

    if FORM_PDF.exists():
        form_doc = fitz.open(FORM_PDF)
        selected_form = [
            (5, "System", "샘플 체계구성도 레이아웃 참고"),
            (10, "Process", "샘플 분석지원 절차 레이아웃 참고"),
            (11, "Process", "샘플 백업/복구 절차 레이아웃 참고"),
            (13, "Reference", "샘플 산출물 관리 표 레이아웃 참고"),
            (14, "Reference", "샘플 수행/관리 표 레이아웃 참고"),
            (15, "Schedule", "샘플 일정/계획 표 레이아웃 참고"),
        ]
        for page_no, category, title in selected_form:
            if page_no > form_doc.page_count:
                continue
            page = form_doc[page_no - 1]
            pix = page.get_pixmap(matrix=fitz.Matrix(1.7, 1.7), alpha=False)
            img_path = ASSETS / f"form_pdf_p{page_no:02d}_{category}.png"
            pix.save(str(img_path))
            gray = ImageOps.grayscale(Image.frombytes("RGB", [pix.width, pix.height], pix.samples))
            score = ImageStat.Stat(gray.filter(ImageFilter.FIND_EDGES)).mean[0]
            candidates.append(
                ImageCandidate(
                    page=page_no,
                    category=category,
                    title=title,
                    path=img_path,
                    score=float(score),
                    image_count=len(page.get_images(full=True)),
                    drawing_count=len(page.get_drawings()),
                    note="form 샘플에서 레이아웃/표현 방식 참고용으로 추출",
                )
            )

    if FORM_WORK_PPTX.exists():
        media_dir = ASSETS / "form_ppt_media"
        media_dir.mkdir(exist_ok=True)
        import zipfile

        with zipfile.ZipFile(FORM_WORK_PPTX) as zf:
            media_names = [n for n in zf.namelist() if n.startswith("ppt/media/")]
            for idx, name in enumerate(media_names[:12], 1):
                suffix = Path(name).suffix or ".png"
                out = media_dir / f"form_ppt_media_{idx:02d}{suffix}"
                out.write_bytes(zf.read(name))
                try:
                    with Image.open(out) as img:
                        gray = ImageOps.grayscale(img.convert("RGB"))
                        score = ImageStat.Stat(gray.filter(ImageFilter.FIND_EDGES)).mean[0]
                except Exception:
                    score = 0.0
                candidates.append(
                    ImageCandidate(
                        page=idx,
                        category="Reference",
                        title="form PPT 내장 이미지 후보",
                        path=out,
                        score=float(score),
                        image_count=1,
                        drawing_count=0,
                        note="사업 내용과 직접 관련 없는 경우 PPT 본문에는 미사용",
                    )
                )

    return candidates


def slide_plan() -> List[SlidePlan]:
    plans = [
        SlidePlan(1, "해군교전분석모델 Build-II 결과분석 제안서", "표지", "Build-II 결과분석 고도화 중심의 제안 범위 제시", "표지", ["결과분석 중심 제안", "RFP 요구사항 대응", "단계적 구현 및 검증 접근"]),
        SlidePlan(2, "제안 개요", "제안 방향 요약", "기능 나열이 아닌 요구사항 충족 구조로 제안", "요약 도식", ["요구사항 기반 설계", "3D 가시화와 자동 재생 고도화", "확률/전자전 효과 통합 분석", "검증 가능한 개발 체계"]),
        SlidePlan(3, "사업 이해도", "사업 맥락 정리", "Build-II는 모의결과 해석성, 재현성, 검증성 강화가 핵심", "RFP 이미지", ["결과분석 업무 중심 체계개발", "전장/플랫폼/무기체계 상호작용 전시", "시험평가 대응 필요"], "System"),
        SlidePlan(4, "RFP 핵심 요구사항 분석", "요구사항 묶음 제시", "결과분석 구조와 구성 요구사항을 추적 가능한 단위로 분해", "요구사항 매트릭스", ["3D 가시화 및 자동 재생", "멀티뷰 및 수중환경 표현", "탐지·명중·기만·전자전 효과 분석", "시험평가 대응"], "Reference"),
        SlidePlan(5, "결과분석 추진전략", "추진 원칙 제시", "데이터-모델-가시화-검증을 단일 흐름으로 연결", "전략 도식", ["로그 기반 재현성 확보", "분석모델 산출 근거 표시", "3D/차트 동기화", "요구사항별 검증 기준 관리"]),
        SlidePlan(6, "To-Be 체계 개념도", "목표 구조 제시", "모의결과를 결과분석 플랫폼에서 통합 가시화", "체계 개념도", ["모의 로그 수집", "분석모델 계산", "3D/멀티뷰 재생", "시험평가 증적 관리"], "Architecture"),
        SlidePlan(7, "결과분석 체계 구조", "기능 블록 정의", "결과분석 기능을 데이터 처리, 분석모델, 전시, 검증으로 모듈화", "구조도", ["데이터 수집·정제", "확률/효과 분석", "3D 가시화·재생", "결과 저장·리포팅"], "Analysis"),
        SlidePlan(8, "3D 가시화 개선방안", "3D 요구사항 대응", "전장상황, 플랫폼, 무기체계, 결과분석 데이터를 동일 장면에서 표현", "3D 도식", ["객체 계층 및 속성 조회", "전장상황 시간축 동기화", "분석 결과 오버레이", "화면 캡처 및 보고 연계"], "3D"),
        SlidePlan(9, "자동 3D 재생 방안", "자동 재생 구조 제시", "시뮬레이션 로그를 이벤트 타임라인으로 재구성", "재생 흐름도", ["로그 파싱 및 시간 정렬", "이벤트 기반 상태 복원", "배속·일시정지·구간 반복", "재현성 검증 로그 생성"], "Simulation"),
        SlidePlan(10, "멀티뷰 기반 전장상황 전시", "멀티뷰 운용방안 제시", "전장환경과 플랫폼/무기체계 뷰를 연동해 분석 효율 향상", "멀티뷰 와이어프레임", ["전장환경 전체 뷰", "함정·항공기·잠수함 추적 뷰", "무기체계 이벤트 뷰", "선택 객체 동기화"], "Analysis"),
        SlidePlan(11, "수중환경 3D 가시화", "수중환경 요구 대응", "수온, 분포, 전달손실, 음선경로를 분석 레이어로 제공", "환경 레이어 도식", ["수온/수중분포 레이어", "음파전달손실 표시", "음선경로 3D 표현", "어뢰·기만 분석과 연계"], "Environment"),
        SlidePlan(12, "계산모델 가시화", "모델 해석성 확보", "입력값-계산단계-출력값을 추적 가능한 화면으로 구성", "모델 흐름도", ["모델 구조 표시", "입력 파라미터 조회", "중간 산출값 확인", "결과 근거 리포트화"]),
        SlidePlan(13, "탐지확률 실시간 분석", "확률 분석 대응", "탐지 이벤트와 환경/센서 조건을 시간축 차트로 표현", "차트 도식", ["탐지확률 추이", "임계값 및 이벤트 표시", "센서/환경 조건 필터", "3D 장면과 동기화"]),
        SlidePlan(14, "명중확률 실시간 분석", "명중 분석 대응", "무기체계 상태와 교전 조건을 반영한 명중확률 분석", "차트 도식", ["명중확률 추이", "교전 이벤트 연동", "입력 조건 비교", "검증 시나리오별 리포트"]),
        SlidePlan(15, "유도탄 기만확률 개선", "유도탄 분석 대응", "기만기 운용 조건과 유도탄 상태를 반영한 성공확률 산출", "분석 흐름도", ["기만 이벤트 식별", "표적/유도탄 상태 반영", "기만 성공확률 산출", "민감도 분석 지원"]),
        SlidePlan(16, "어뢰 기만확률 개선", "어뢰 분석 대응", "수중환경과 음향 조건을 포함해 어뢰 기만 효과를 분석", "분석 흐름도", ["수중환경 조건 입력", "어뢰 탐지/추적 상태 반영", "기만기 효과 산출", "환경별 결과 비교"]),
        SlidePlan(17, "전자전 효과 분석", "전자전 요구 대응", "재머, GPS 교란 등 전자전 효과를 교전결과와 연계 분석", "전자전 흐름도", ["전자전 이벤트 수집", "재머 효과 산출", "GPS 교란 영향 분석", "탐지·명중 결과 연계"]),
        SlidePlan(18, "데이터 처리 흐름", "데이터 흐름 정의", "모의 로그부터 분석 리포트까지 추적 가능한 파이프라인 구성", "데이터 흐름도", ["로그 수집", "정합성 검증", "분석모델 계산", "3D/차트 전시", "결과 저장"]),
        SlidePlan(19, "시스템 아키텍처", "시스템 구조 제시", "표준 인터페이스와 모듈화 구조로 확장성 확보", "아키텍처 도식", ["데이터 계층", "분석 서비스 계층", "가시화 계층", "검증/운영 지원 계층"], "Architecture"),
        SlidePlan(20, "UI/UX 설계 방향", "운용 화면 방향 제시", "분석 담당자가 빠르게 상황을 재현하고 근거를 확인하도록 구성", "화면 구성", ["시나리오 선택", "3D 재생·멀티뷰", "확률/효과 차트", "결과 리포트"]),
        SlidePlan(21, "개발 방법론", "수행방식 제시", "요구사항 추적성과 산출물 검증을 중심으로 개발 수행", "프로세스", ["요구사항 분석", "반복 설계/구현", "통합 시험", "검증 증적 관리"], "Process"),
        SlidePlan(24, "시험평가 대응", "시험평가 방안 제시", "요구사항별 시험항목과 판정기준을 사전에 정의", "시험 매트릭스", ["요구사항-시험항목 매핑", "시나리오 기반 기능시험", "성능/안정성 확인", "결함 조치 추적"], "Schedule"),
        SlidePlan(26, "품질보증 계획", "품질 관리 제시", "산출물, 코드, 모델, 시험결과를 기준에 따라 관리", "QA 체크리스트", ["품질 기준 수립", "리뷰 및 승인 절차", "결함 관리", "변경 영향 분석"]),
        SlidePlan(27, "일정 계획", "단계별 일정 제시", "분석-설계-구현-통합-시험평가를 단계적으로 수행", "일정표", ["착수 및 요구사항 분석", "핵심 기능 설계", "반복 구현 및 통합", "시험평가 및 안정화"], "Schedule"),
        SlidePlan(28, "투입 인력 계획", "역할 중심 계획", "확인 가능한 범위에서 역할 기반 투입 구조를 제시", "역할 표", ["사업관리", "결과분석/모델", "3D 가시화/UI", "시험평가/V&V"], "Organization"),
        SlidePlan(29, "기대효과", "사업효과 정리", "분석 효율, 재현성, 검증성을 중심으로 기대효과 제시", "효과 도식", ["모의결과 해석성 강화", "자동 재생 기반 분석 효율 향상", "확률/전자전 효과 통합 판단", "확장 가능한 분석 플랫폼 기반"]),
        SlidePlan(30, "결론", "최종 메시지", "Build-II 요구사항에 부합하는 단계적 구현 및 검증 중심 접근", "종합 메시지", ["결과분석 고도화를 통해 모의결과의 해석성, 재현성, 검증성을 강화", "3D 전장 가시화와 자동 재생을 통해 분석 효율 향상", "탐지·명중·기만·전자전 효과를 통합 분석하여 의사결정 지원 능력 강화", "Build-II 체계개발 요구사항에 부합하는 단계적 구현 및 검증 중심 접근", "향후 확장 가능한 결과분석 플랫폼 기반 확보"]),
    ]
    removed_titles = {"투입 인력 계획"}
    return [p for p in plans if p.title not in removed_titles]


def write_markdown(task_text: str, images: List[ImageCandidate], plans: List[SlidePlan]) -> None:
    reqs = requirements()
    (OUTPUT / "proposal_outline.md").write_text(
        "\n".join(
            [
                "# 해군교전분석모델 Build-II 결과분석 제안서 개요",
                "",
                "## 제안 방향",
                "- 기능 소개 중심이 아니라 RFP 요구사항 대응 중심으로 구성한다.",
                "- 결과분석 데이터, 3D 가시화, 자동 재생, 확률/전자전 효과 분석, 시험평가 대응을 핵심 축으로 둔다.",
                "- 확인 불가한 회사명, 실적, 인력, 인증은 기재하지 않고 역할·제안 방향으로 표기한다.",
                "",
                "## Task 반영 요약",
                task_text.strip(),
                "",
                "## 장표 구성",
                f"- 총 {len(plans)}장 구성",
                "- 한 장표당 Bullet 3~5개 기준",
                "- RFP 원문 이미지와 PowerPoint 도형 기반 자동 도식을 병행 사용",
            ]
        ),
        encoding="utf-8",
    )

    lines = ["# 추출 요구사항", "", "RFP PDF는 텍스트 레이어가 없는 스캔형 문서로 확인되어, 사용자 제공 티켓과 task.txt를 기준 요구사항으로 구조화하고 PDF 페이지 이미지를 근거 자료로 분류하였다.", ""]
    lines += ["| 구분 | 요구사항 | 대응 방향 | 검증 방안 | 우선순위 |", "|---|---|---|---|---|"]
    for r in reqs:
        lines.append(f"| {r.area} | {r.item} | {r.response} | {r.verification} | {r.priority} |")
    (OUTPUT / "extracted_requirements.md").write_text("\n".join(lines), encoding="utf-8")

    lines = ["# 슬라이드 계획", "", "| No. | 제목 | 목적 | 주요 메시지 | 시각자료 |", "|---:|---|---|---|---|"]
    for p in plans:
        lines.append(f"| {p.no} | {p.title} | {p.purpose} | {p.message} | {p.visual} |")
    (OUTPUT / "slide_plan.md").write_text("\n".join(lines), encoding="utf-8")

    lines = [
        "# 이미지 카탈로그",
        "",
        "- `해군교전2요청서.pdf`는 본 프로젝트 핵심 내용으로 전체 85페이지를 분석했으며, PPT 본문에는 원문 페이지 이미지를 사용하지 않고 요구사항 근거로만 반영하였다.",
        "- `pdf_page_analysis.csv`에는 RFP 페이지별 텍스트 레이어 길이, 이미지 수, 도형 수, 선명도/정보밀도 지표를 기록하였다.",
        "- `form` 폴더의 PDF/PPT에서 레이아웃과 도식 스타일 참고 후보를 추출하였다.",
        "- 기존 샘플의 타 사업명/유지관리 문구가 본 사업 내용과 혼재되지 않도록, 최종 PPT 본문은 신규 표·도형 중심으로 구성하였다.",
        "",
        "| Page | Category | Title | File | Quality Score | Note |",
        "|---:|---|---|---|---:|---|",
    ]
    for img in images:
        rel = img.path.relative_to(OUTPUT).as_posix()
        lines.append(f"| {img.page} | {img.category} | {img.title} | `{rel}` | {img.score:.2f} | {img.note} |")
    (OUTPUT / "image_catalog.md").write_text("\n".join(lines), encoding="utf-8")

    (OUTPUT / "README.md").write_text(
        "\n".join(
            [
                "# 제안서 자동 생성 산출물",
                "",
                "## 실행 방법",
                "```powershell",
                "py build_ppt.py",
                "```",
                "",
                "계획 문서만 갱신하려면 다음을 실행한다.",
                "",
                "```powershell",
                "py build_ppt.py --plan-only",
                "```",
                "",
                "## 주요 산출물",
                "- `proposal_outline.md`: 제안서 개요",
                "- `slide_plan.md`: 슬라이드별 목적 및 메시지",
                "- `extracted_requirements.md`: 요구사항 대응 매트릭스",
                "- `image_catalog.md`: PDF 이미지 후보 분류",
                "- `assets/`: PDF 렌더링 이미지와 접촉시트",
                "- `해군교전분석모델_BuildII_결과분석_제안서_초안.pptx`: 제안서 PPT",
            ]
        ),
        encoding="utf-8",
    )


def find_image(images: List[ImageCandidate], category: Optional[str]) -> Optional[Path]:
    if not category:
        return None
    for img in images:
        if img.category == category:
            return img.path
    return None


def set_text(shape, text: str, size: int = 11, bold: bool = False, color: RGBColor = DARK, align=None, font: str = "Malgun Gothic") -> None:
    shape.text_frame.clear()
    shape.text_frame.word_wrap = True
    shape.text_frame.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
    shape.text_frame.margin_left = Cm(0.1)
    shape.text_frame.margin_right = Cm(0.1)
    shape.text_frame.margin_top = Cm(0.05)
    shape.text_frame.margin_bottom = Cm(0.05)
    p = shape.text_frame.paragraphs[0]
    p.text = text
    p.font.name = font
    p.font.size = Pt(size)
    p.font.bold = bold
    p.font.color.rgb = color
    if align:
        p.alignment = align


def add_box(slide, x, y, w, h, text, fill=WHITE, line=LINE, size=10, bold=False, color=DARK, font: str = "Malgun Gothic"):
    box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, int(x), int(y), int(w), int(h))
    box.fill.solid()
    box.fill.fore_color.rgb = fill
    box.line.color.rgb = line
    set_text(box, text, size=size, bold=bold, color=color, align=PP_ALIGN.CENTER, font=font)
    box.text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
    return box


def add_arrow(slide, x, y, w, h):
    sh = slide.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, int(x), int(y), int(w), int(h))
    sh.fill.solid()
    sh.fill.fore_color.rgb = SKY
    sh.line.color.rgb = SKY
    return sh


def apply_custom_background(slide, width, height, cover: bool = False):
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = WHITE


def add_header(slide, plan: SlidePlan):
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = WHITE
    band = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Cm(19.05), Cm(1.15))
    band.fill.solid()
    band.fill.fore_color.rgb = NAVY
    band.line.color.rgb = NAVY
    title = slide.shapes.add_textbox(Cm(1.2), Cm(0.23), Cm(15.2), Cm(0.62))
    set_text(title, plan.title, size=15, bold=True, color=WHITE)
    no = slide.shapes.add_textbox(Cm(17.0), Cm(0.25), Cm(1.5), Cm(0.5))
    set_text(no, f"{plan.no:02d}", size=12, bold=True, color=SKY, align=PP_ALIGN.RIGHT)
    rule = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Cm(0.65), Cm(1.45), Cm(17.75), Cm(0.05))
    rule.fill.solid()
    rule.fill.fore_color.rgb = BLUE
    rule.line.color.rgb = BLUE
    msg = slide.shapes.add_textbox(Cm(0.75), Cm(1.65), Cm(17.55), Cm(0.8))
    set_text(msg, plan.message, size=11, bold=True, color=NAVY)


def add_footer(slide):
    foot = slide.shapes.add_textbox(Cm(0.75), Cm(26.5), Cm(17.6), Cm(0.4))
    set_text(foot, "해군교전분석모델 Build-II 결과분석 제안서", size=7, color=GRAY)


def add_bullets(slide, bullets: Iterable[str], x=Cm(1.0), y=Cm(20.6), w=Cm(17.0), h=Cm(4.2)):
    box = slide.shapes.add_textbox(x, y, w, h)
    tf = box.text_frame
    tf.clear()
    tf.margin_left = Cm(0.1)
    for i, b in enumerate(bullets):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = b
        p.level = 0
        p.font.name = "Malgun Gothic"
        p.font.size = Pt(10)
        p.font.color.rgb = DARK
        p.space_after = Pt(3)


def add_image_area(slide, image_path: Optional[Path], fallback_title: str):
    if image_path and image_path.exists():
        slide.shapes.add_picture(str(image_path), Cm(1.0), Cm(3.0), width=Cm(17.0), height=Cm(16.7))
        cap = slide.shapes.add_textbox(Cm(1.0), Cm(19.8), Cm(17.0), Cm(0.45))
        set_text(cap, f"RFP 원문 이미지 활용: {image_path.name}", size=7, color=GRAY, align=PP_ALIGN.RIGHT)
    else:
        add_box(slide, Cm(1.2), Cm(3.2), Cm(16.5), Cm(3.0), fallback_title, fill=LIGHT, size=14, bold=True, color=NAVY)


def add_reference_image(slide, image_path: Optional[Path], x=Cm(12.2), y=Cm(11.2), w=Cm(5.4), h=Cm(7.6)):
    if not image_path or not image_path.exists():
        return
    frame = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x - Cm(0.08), y - Cm(0.08), w + Cm(0.16), h + Cm(0.45))
    frame.fill.solid()
    frame.fill.fore_color.rgb = WHITE
    frame.line.color.rgb = LINE
    slide.shapes.add_picture(str(image_path), x, y, width=w, height=h)
    cap = slide.shapes.add_textbox(x, y + h + Cm(0.08), w, Cm(0.35))
    set_text(cap, f"RFP p.{image_path.stem.split('_')[1][1:]}", size=6, color=GRAY, align=PP_ALIGN.RIGHT)


def add_flow(slide, labels: List[str], y=Cm(7.0)):
    x = Cm(1.0)
    for i, label in enumerate(labels):
        add_box(slide, x, y, Cm(3.2), Cm(1.8), label, fill=LIGHT, line=BLUE, size=9, bold=True, color=NAVY)
        x += Cm(3.45)
        if i < len(labels) - 1:
            add_arrow(slide, x, y + Cm(0.55), Cm(0.65), Cm(0.65))
            x += Cm(0.8)


def add_matrix(slide):
    headers = ["요구사항", "대응 기능", "검증 관점"]
    rows = [
        ["3D/자동재생", "로그 기반 3D 재현", "시간축·이벤트 재현성"],
        ["멀티뷰/수중환경", "동기화 뷰·환경 레이어", "입력자료-화면 대조"],
        ["확률/전자전", "실시간 차트·효과 산출", "산출식·시나리오 검증"],
        ["시험평가", "요구사항별 검증 및 산출물 확인", "시험 결과·검증 증적"],
    ]
    x0, y0 = Cm(1.0), Cm(4.0)
    widths = [Cm(4.6), Cm(6.0), Cm(6.0)]
    for j, h in enumerate(headers):
        add_box(slide, x0 + sum(widths[:j]), y0, widths[j], Cm(1.0), h, fill=NAVY, line=NAVY, size=9, bold=True, color=WHITE)
    for i, row in enumerate(rows):
        for j, cell in enumerate(row):
            add_box(slide, x0 + sum(widths[:j]), y0 + Cm(1.0 + i * 1.4), widths[j], Cm(1.25), cell, fill=WHITE if i % 2 else LIGHT, line=LINE, size=8, color=DARK)


def add_chart_mock(slide, title: str):
    add_box(slide, Cm(1.2), Cm(4.0), Cm(16.5), Cm(12.0), "", fill=WHITE, line=LINE)
    tx = slide.shapes.add_textbox(Cm(1.6), Cm(4.25), Cm(10), Cm(0.5))
    set_text(tx, title, size=11, bold=True, color=NAVY)
    left, top, width, height = Cm(2.0), Cm(6.0), Cm(14.2), Cm(8.0)
    for i in range(5):
        y = top + height * i / 4
        line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, y, width, Cm(0.02))
        line.fill.solid(); line.fill.fore_color.rgb = LINE; line.line.color.rgb = LINE
    points = [(0, 0.75), (0.16, 0.55), (0.32, 0.62), (0.48, 0.35), (0.64, 0.42), (0.82, 0.25), (1.0, 0.30)]
    last = None
    for px, py in points:
        x = left + width * px
        y = top + height * py
        dot = slide.shapes.add_shape(MSO_SHAPE.OVAL, x - Cm(0.08), y - Cm(0.08), Cm(0.16), Cm(0.16))
        dot.fill.solid(); dot.fill.fore_color.rgb = BLUE; dot.line.color.rgb = BLUE
        if last:
            lx, ly = last
            seg = slide.shapes.add_connector(1, lx, ly, x, y)
            seg.line.color.rgb = BLUE
            seg.line.width = Pt(2)
        last = (x, y)
    add_box(slide, Cm(2.0), Cm(14.5), Cm(3.2), Cm(0.7), "이벤트", fill=SKY, line=SKY, size=8, bold=True, color=NAVY)
    add_box(slide, Cm(5.5), Cm(14.5), Cm(3.2), Cm(0.7), "임계값", fill=LIGHT, line=LINE, size=8)
    add_box(slide, Cm(9.0), Cm(14.5), Cm(3.2), Cm(0.7), "산출근거", fill=LIGHT, line=LINE, size=8)


def add_architecture(slide):
    layers = [
        ("가시화 계층", ["3D 전장상황", "멀티뷰", "실시간 차트"]),
        ("분석 서비스 계층", ["탐지/명중", "기만확률", "전자전 효과"]),
        ("데이터 처리 계층", ["로그 수집", "정합성 검증", "시간축 정렬"]),
        ("검증/운영 계층", ["요구사항 추적", "시험 검증", "시험 리포트"]),
    ]
    y = Cm(4.0)
    for name, items in layers:
        add_box(slide, Cm(1.0), y, Cm(3.8), Cm(1.4), name, fill=NAVY, line=NAVY, size=9, bold=True, color=WHITE)
        for i, item in enumerate(items):
            add_box(slide, Cm(5.2 + i * 4.1), y, Cm(3.7), Cm(1.4), item, fill=LIGHT, line=BLUE, size=9, bold=True, color=NAVY)
        y += Cm(2.35)


def add_form_header(slide, plan: SlidePlan, width, height):
    apply_custom_background(slide, width, height)
    chap = slide.shapes.add_textbox(Cm(0.75), Cm(0.48), Cm(2.0), Cm(0.35))
    set_text(chap, f"Ⅱ - {max(plan.no - 1, 1)}", size=7.5, bold=True, color=NAVY, font="Malgun Gothic")
    title = slide.shapes.add_textbox(Cm(0.75), Cm(0.88), width - Cm(1.5), Cm(0.78))
    set_text(title, plan.title, size=11.5, bold=True, color=DARK, font="Malgun Gothic")
    line1 = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Cm(0.75), Cm(1.72), width - Cm(1.5), Cm(0.03))
    line1.fill.solid()
    line1.fill.fore_color.rgb = NAVY
    line1.line.color.rgb = NAVY
    line2 = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Cm(0.75), Cm(1.77), Cm(4.4), Cm(0.08))
    line2.fill.solid()
    line2.fill.fore_color.rgb = BLUE
    line2.line.color.rgb = BLUE
    msg = slide.shapes.add_textbox(Cm(0.95), Cm(2.02), width - Cm(1.9), Cm(0.62))
    set_text(msg, plan.message, size=7.2, bold=True, color=NAVY, font="Malgun Gothic")


def slide_detail(plan: SlidePlan):
    tags = {
        3: ["사업이해", "RFP 분석"],
        4: ["요구사항", "추적성"],
        8: ["3D 가시화", "전장상황"],
        9: ["자동 재생", "로그 기반"],
        10: ["멀티뷰", "상호작용"],
        11: ["수중환경", "음향 분석"],
        12: ["계산모델", "해석성"],
        13: ["탐지확률", "실시간 차트"],
        14: ["명중확률", "교전 분석"],
        15: ["유도탄", "기만확률"],
        16: ["어뢰", "기만확률"],
        17: ["전자전", "재머/GPS"],
        21: ["개발방법론", "추적성"],
        24: ["시험평가", "판정기준"],
        27: ["일정", "단계관리"],
    }.get(plan.no, [plan.purpose, "제안 방향"])
    if plan.no in CORE_SLIDE_DETAILS:
        paragraph = CORE_SLIDE_DETAILS[plan.no]["objective"]
    else:
        paragraph = (
            f"{plan.message}. 요구사항 충족 여부를 확인할 수 있도록 입력자료, 처리절차, 산출물, "
            "검증기준을 함께 제시하여 사업 수행 과정에서의 추적성과 검증성을 확보한다."
        )
    return tags, paragraph


def add_form_intro(slide, plan: SlidePlan, width):
    _, paragraph = slide_detail(plan)
    subsection = {
        2: "가. 제안 개요",
        3: "나. 사업 이해도",
        4: "다. 핵심 요구사항 분석",
        5: "라. 결과분석 추진전략",
        6: "마. To-Be 체계 개념",
        7: "바. 결과분석 체계 구조",
        8: "1) 결과분석 데이터 3D 가시화 개선",
        9: "2) 모의결과 기반 자동 3D 재생",
        10: "3) 멀티뷰 기반 전장상황 전시",
        11: "4) 수중환경 3D 가시화",
        12: "5) 계산모델 가시화",
        13: "6) 탐지확률 실시간 분석",
        14: "7) 명중확률 실시간 분석",
        15: "8) 유도탄 기만확률 개선",
        16: "9) 어뢰 기만확률 개선",
        17: "10) 전자전 효과 분석",
        18: "가. 데이터 처리 흐름",
        19: "나. 시스템 아키텍처",
        20: "다. UI/UX 설계 방향",
        21: "라. 개발 방법론",
        24: "마. 시험평가 대응",
        26: "사. 품질보증 계획",
        27: "아. 일정 계획",
        29: "자. 기대효과",
        30: "차. 결론",
    }.get(plan.no, "가. 제안 내용")
    label = slide.shapes.add_textbox(Cm(0.95), Cm(2.75), width - Cm(1.9), Cm(0.36))
    set_text(label, subsection, size=7.5, bold=True, color=NAVY, font="Malgun Gothic")
    body = slide.shapes.add_textbox(Cm(0.95), Cm(3.18), width - Cm(1.9), Cm(0.78))
    set_text(body, paragraph, size=7.0, color=DARK, font="Malgun Gothic")


def add_form_footer(slide, page_no: int, width, height):
    foot = slide.shapes.add_textbox(Cm(0.75), height - Cm(0.75), width - Cm(1.5), Cm(0.35))
    set_text(foot, f"해군교전분석모델 Build-II 결과분석 제안서    {page_no:02d}", size=6, color=GRAY, align=PP_ALIGN.RIGHT)


def add_form_bullets(slide, bullets: Iterable[str], width, height):
    y = height - Cm(4.15)
    title = slide.shapes.add_textbox(Cm(0.85), y, Cm(3.0), Cm(0.4))
    set_text(title, "핵심 포인트", size=8.2, bold=True, color=NAVY)
    box = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, int(Cm(0.75)), int(y + Cm(0.48)), int(width - Cm(1.5)), int(Cm(2.65)))
    box.fill.solid()
    box.fill.fore_color.rgb = PALE_GRAY
    box.line.color.rgb = LINE
    tf = box.text_frame
    tf.clear()
    tf.word_wrap = True
    tf.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
    tf.margin_left = Cm(0.32)
    tf.margin_right = Cm(0.28)
    tf.margin_top = Cm(0.14)
    tf.margin_bottom = Cm(0.1)
    for i, bullet in enumerate(list(bullets)[:4]):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = f"• {bullet}"
        p.alignment = PP_ALIGN.LEFT
        p.level = 0
        p.font.name = "Malgun Gothic"
        p.font.size = Pt(6.9)
        p.font.color.rgb = DARK
        p.space_after = Pt(1)


def add_form_table(slide, headers: List[str], rows: List[List[str]], x, y, widths, row_h=Cm(1.05)):
    for j, h in enumerate(headers):
        add_box(slide, x + sum(widths[:j]), y, widths[j], row_h, h, fill=NAVY, line=NAVY, size=7.0, bold=True, color=WHITE)
    for i, row in enumerate(rows):
        for j, cell in enumerate(row):
            box = add_box(slide, x + sum(widths[:j]), y + row_h * (i + 1), widths[j], row_h, cell, fill=WHITE if i % 2 else PALE_BLUE, line=LINE, size=6.7, color=DARK)
            if j > 0:
                p = box.text_frame.paragraphs[0]
                p.alignment = PP_ALIGN.LEFT
                box.text_frame.margin_left = Cm(0.18)
                box.text_frame.margin_right = Cm(0.18)


def add_form_process(slide, labels: List[str], width, y=Cm(6.3)):
    usable = width - Cm(2.0)
    box_w = min(Cm(3.2), usable / max(len(labels), 1) - Cm(0.35))
    gap = (usable - box_w * len(labels)) / max(len(labels) - 1, 1)
    x = Cm(1.0)
    for i, label in enumerate(labels):
        add_box(slide, x, y, box_w, Cm(2.2), label, fill=WHITE, line=BLUE, size=7.4, bold=True, color=NAVY)
        if i < len(labels) - 1:
            add_arrow(slide, x + box_w + gap / 2 - Cm(0.25), y + Cm(0.8), Cm(0.5), Cm(0.5))
        x += box_w + gap


def add_form_cards(slide, labels: List[str], width, y=Cm(5.35)):
    cols = 2 if len(labels) <= 4 else 3
    card_w = (width - Cm(2.0) - Cm(0.45) * (cols - 1)) / cols
    card_h = Cm(2.6)
    for i, label in enumerate(labels):
        row, col = divmod(i, cols)
        x = Cm(1.0) + col * (card_w + Cm(0.45))
        cy = y + row * (card_h + Cm(0.5))
        add_box(slide, x, cy, card_w, card_h, label, fill=PALE_BLUE if i % 2 == 0 else WHITE, line=BLUE, size=8.3, bold=True, color=NAVY)


def add_form_architecture(slide, width):
    layers = [
        ("가시화", ["3D 전장", "멀티뷰", "차트"]),
        ("분석모델", ["탐지", "명중", "기만", "전자전"]),
        ("데이터", ["로그", "이벤트", "환경자료"]),
        ("검증", ["시험", "추적성", "증적"]),
    ]
    y = Cm(4.85)
    for idx, (layer, items) in enumerate(layers):
        add_box(slide, Cm(0.95), y, Cm(3.0), Cm(1.2), layer, fill=NAVY if idx == 0 else BLUE, line=BLUE, size=8.3, bold=True, color=WHITE)
        item_w = (width - Cm(5.2)) / len(items)
        for i, item in enumerate(items):
            add_box(slide, Cm(4.25) + item_w * i, y, item_w - Cm(0.2), Cm(1.2), item, fill=WHITE, line=LINE, size=7.6, bold=True, color=DARK)
        y += Cm(1.55)


def add_form_chart(slide, width, title: str):
    x, y, w, h = Cm(1.0), Cm(5.35), width - Cm(2.0), Cm(6.7)
    add_box(slide, x, y, w, h, "", fill=WHITE, line=LINE)
    label = slide.shapes.add_textbox(x + Cm(0.35), y + Cm(0.2), w - Cm(0.7), Cm(0.4))
    set_text(label, title, size=8.5, bold=True, color=NAVY)
    left, top, cw, ch = x + Cm(0.8), y + Cm(1.4), w - Cm(1.5), h - Cm(2.2)
    for i in range(5):
        grid = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, int(left), int(top + ch * i / 4), int(cw), int(Cm(0.015)))
        grid.fill.solid()
        grid.fill.fore_color.rgb = LINE
        grid.line.color.rgb = LINE
    series = [(0.0, 0.68), (0.15, 0.55), (0.31, 0.60), (0.46, 0.42), (0.64, 0.35), (0.82, 0.28), (1.0, 0.32)]
    last = None
    for px, py in series:
        sx, sy = left + cw * px, top + ch * py
        dot = slide.shapes.add_shape(MSO_SHAPE.OVAL, int(sx - Cm(0.06)), int(sy - Cm(0.06)), int(Cm(0.12)), int(Cm(0.12)))
        dot.fill.solid()
        dot.fill.fore_color.rgb = BLUE
        dot.line.color.rgb = BLUE
        if last:
            conn = slide.shapes.add_connector(1, int(last[0]), int(last[1]), int(sx), int(sy))
            conn.line.color.rgb = BLUE
            conn.line.width = Pt(1.7)
        last = (sx, sy)


def add_form_timeline(slide, width):
    rows = [
        ["착수/분석", "요구사항 정제, RFP 대응 매트릭스 수립"],
        ["설계", "3D 가시화, 분석모델, 데이터 흐름 상세설계"],
        ["구현", "스프린트 기반 기능 구현 및 통합"],
        ["시험평가", "요구사항별 검증, 시험 증적 정리"],
    ]
    add_form_table(slide, ["단계", "주요 수행내용"], rows, Cm(1.0), Cm(4.85), [Cm(3.5), width - Cm(5.5)], row_h=Cm(1.08))


CORE_SLIDE_DETAILS = {
    8: {
        "objective": "전장상황, 플랫폼, 무기체계, 결과분석 데이터를 하나의 3D 분석 장면에서 통합 표현",
        "steps": ["전장 객체\n표준화", "시간/좌표\n동기화", "분석결과\n오버레이", "분석장면\n저장"],
        "table": [
            ["구현 기능", "플랫폼/무기체계 궤적, 교전 이벤트, 결과분석 값을 3D 객체 속성으로 연계"],
            ["전시 방식", "객체 색상·라벨·궤적·이벤트 아이콘으로 상태와 결과를 구분"],
            ["검증 관점", "동일 로그 재생 시 위치, 시간, 이벤트 표시가 반복 재현되는지 확인"],
        ],
        "mockup": "3d",
    },
    9: {
        "objective": "모의 로그와 이벤트를 시간축으로 재구성하여 자동 3D 재생과 구간 분석 지원",
        "steps": ["로그 수집", "이벤트\n정렬", "상태 복원", "자동 재생"],
        "table": [
            ["구현 기능", "재생, 일시정지, 배속, 구간 반복, 이벤트 점프, 장면 북마크 제공"],
            ["전시 방식", "3D 장면과 이벤트 타임라인을 동기화하여 분석자가 원인 구간을 추적"],
            ["검증 관점", "이벤트 발생 시각과 화면 상태가 로그 기준값과 일치하는지 확인"],
        ],
        "mockup": "timeline",
    },
    10: {
        "objective": "전장환경, 함정, 항공기, 잠수함, 무기체계를 동시에 관찰하는 멀티뷰 분석 화면 구성",
        "steps": ["전체 전장", "함정/항공기", "잠수함", "무기체계"],
        "table": [
            ["구현 기능", "뷰 분할, 선택 객체 동기화, 필터, 확대/추적, 이벤트 연동"],
            ["전시 방식", "전체 상황과 개별 플랫폼/무기체계 상세 뷰를 동시에 제공"],
            ["검증 관점", "한 뷰에서 선택한 객체가 다른 뷰와 차트에 동일하게 반영되는지 확인"],
        ],
        "mockup": "multiview",
    },
    11: {
        "objective": "수온, 수중분포, 음파전달손실, 음선경로를 3D 환경 레이어로 제공",
        "steps": ["환경자료\n입력", "수중 레이어\n생성", "음향 경로\n표현", "무기분석\n연계"],
        "table": [
            ["구현 기능", "수온/분포 단면, 전달손실 등고면, 음선경로 곡선, 시간별 환경 변화 표시"],
            ["전시 방식", "수중환경 레이어를 어뢰·기만 분석 화면과 선택적으로 중첩"],
            ["검증 관점", "입력 환경자료와 시각화 레이어 값, 범례, 단위가 일치하는지 확인"],
        ],
        "mockup": "underwater",
    },
    12: {
        "objective": "계산모델의 입력값, 계산 흐름, 중간 산출값, 출력값을 추적 가능한 형태로 가시화",
        "steps": ["입력값", "모델 구조", "중간 산출", "출력/근거"],
        "table": [
            ["구현 기능", "모델별 입력 파라미터, 계산 단계, 산출식 버전, 출력값 조회"],
            ["전시 방식", "계산 흐름도와 입력/출력 표를 연결하여 산출 근거를 표시"],
            ["검증 관점", "산출값이 로그, 모델 버전, 입력 파라미터와 추적되는지 확인"],
        ],
        "mockup": "model",
    },
    13: {
        "objective": "탐지 이벤트, 센서 조건, 환경 조건에 따른 탐지확률 변화를 실시간 차트로 표현",
        "steps": ["센서/환경\n입력", "탐지 이벤트", "확률 산출", "차트/3D\n동기화"],
        "table": [
            ["구현 기능", "탐지확률 시계열, 임계값, 탐지 이벤트 마커, 조건 필터 제공"],
            ["전시 방식", "3D 객체 선택 시 해당 센서/표적의 탐지확률 차트를 자동 갱신"],
            ["검증 관점", "동일 입력 조건에서 확률 산출과 차트 표시가 일관되는지 확인"],
        ],
        "mockup": "chart",
    },
    14: {
        "objective": "무기체계 상태, 거리, 교전 조건을 반영하여 명중확률을 시간축으로 분석",
        "steps": ["교전 조건", "무기 상태", "명중확률\n산출", "결과 비교"],
        "table": [
            ["구현 기능", "교전 구간별 명중확률, 사격/유도 이벤트, 조건별 비교 차트 제공"],
            ["전시 방식", "무기 발사-유도-명중/실패 흐름과 확률 변화를 함께 표시"],
            ["검증 관점", "시나리오별 명중 이벤트와 확률 산출 결과의 정합성 확인"],
        ],
        "mockup": "chart",
    },
    15: {
        "objective": "유도탄과 기만기 상호작용을 반영하여 기만 성공확률 산출 절차 개선",
        "steps": ["기만 이벤트", "유도탄 상태", "기만 효과", "성공확률"],
        "table": [
            ["구현 기능", "기만기 투발 시점, 유도탄 탐색/추적 상태, 표적 거리·방위 조건 반영"],
            ["전시 방식", "기만 전후 유도탄 상태 변화와 성공확률을 타임라인으로 표시"],
            ["검증 관점", "대표 시나리오 회귀시험과 입력 조건별 민감도 비교 수행"],
        ],
        "mockup": "decoy",
    },
    16: {
        "objective": "수중환경과 음향 조건을 반영하여 어뢰 기만 성공확률 산출 개선",
        "steps": ["수중환경", "어뢰 탐지", "기만기 효과", "확률 비교"],
        "table": [
            ["구현 기능", "음파전달손실, 음선경로, 어뢰 탐지/추적 상태, 기만기 효과 반영"],
            ["전시 방식", "수중환경 레이어와 어뢰/기만기 궤적, 성공확률 차트를 연계"],
            ["검증 관점", "환경 조건 변화에 따른 산출값 차이를 기준 시나리오로 비교"],
        ],
        "mockup": "underwater",
    },
    17: {
        "objective": "재머, GPS 교란 등 전자전 장비 효과를 탐지·명중·기만 결과와 연계 분석",
        "steps": ["전자전 이벤트", "재머 효과", "GPS 교란", "교전 영향"],
        "table": [
            ["구현 기능", "재머 영향권, GPS 교란 이벤트, 전자전 효과 지표, 영향 대상 추적"],
            ["전시 방식", "전자전 영향 범위와 확률/교전 결과 변화를 한 화면에서 표시"],
            ["검증 관점", "전자전 이벤트 발생 전후 탐지·명중 결과 변화와 효과 산출값 대조"],
        ],
        "mockup": "ew",
    },
}


def add_section_label(slide, text, x, y, w=Cm(4.0)):
    label = slide.shapes.add_textbox(x, y, w, Cm(0.35))
    set_text(label, text, size=8, bold=True, color=NAVY)


def add_mini_label(slide, text, x, y, w=Cm(2.0), fill=PALE_BLUE):
    return add_box(slide, x, y, w, Cm(0.42), text, fill=fill, line=BLUE, size=6.2, bold=True, color=NAVY)


def add_core_mockup(slide, kind: str, x, y, w, h):
    add_box(slide, x, y, w, h, "", fill=WHITE, line=LINE)
    if kind in ("3d", "decoy", "ew"):
        for i in range(5):
            gx = x + Cm(0.55) + (w - Cm(1.1)) * i / 4
            grid = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, int(gx), int(y + Cm(0.45)), int(Cm(0.012)), int(h - Cm(0.9)))
            grid.fill.solid(); grid.fill.fore_color.rgb = LINE; grid.line.color.rgb = LINE
        sea = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, int(x + Cm(0.25)), int(y + h - Cm(1.05)), int(w - Cm(0.5)), int(Cm(0.35)))
        sea.fill.solid(); sea.fill.fore_color.rgb = PALE_BLUE; sea.line.color.rgb = PALE_BLUE
        add_box(slide, x + Cm(0.6), y + Cm(0.8), Cm(1.9), Cm(0.62), "아군 함정", fill=PALE_BLUE, line=BLUE, size=6.5, bold=True)
        add_box(slide, x + w - Cm(2.7), y + Cm(1.75), Cm(1.9), Cm(0.62), "표적", fill=LIGHT, line=LINE, size=6.5, bold=True)
        conn = slide.shapes.add_connector(1, int(x + Cm(2.5)), int(y + Cm(1.1)), int(x + w - Cm(2.7)), int(y + Cm(2.05)))
        conn.line.color.rgb = BLUE; conn.line.width = Pt(1.4)
        add_mini_label(slide, "궤적/이벤트", x + Cm(3.0), y + Cm(0.72), Cm(2.35))
        add_mini_label(slide, "분석값 오버레이", x + w - Cm(4.15), y + Cm(0.72), Cm(3.0), fill=WHITE)
        if kind == "decoy":
            add_box(slide, x + Cm(3.15), y + Cm(2.35), Cm(1.65), Cm(0.55), "기만기", fill=SKY, line=BLUE, size=6.2, bold=True)
            add_mini_label(slide, "성공확률", x + Cm(5.05), y + Cm(2.4), Cm(2.0), fill=WHITE)
        if kind == "ew":
            jammer = slide.shapes.add_shape(MSO_SHAPE.OVAL, int(x + Cm(3.0)), int(y + Cm(1.45)), int(Cm(2.1)), int(Cm(2.1)))
            jammer.fill.solid(); jammer.fill.fore_color.rgb = PALE_BLUE; jammer.line.color.rgb = BLUE
            set_text(jammer, "재머\n영향권", size=6.5, bold=True, color=NAVY, align=PP_ALIGN.CENTER)
            add_mini_label(slide, "GPS 교란 이벤트", x + w - Cm(4.1), y + Cm(2.75), Cm(3.1), fill=WHITE)
    elif kind == "multiview":
        boxes = [("전장환경", 0, 0, 0.58, 0.55), ("플랫폼", 0.62, 0, 0.38, 0.26), ("무기체계", 0.62, 0.31, 0.38, 0.24), ("차트", 0, 0.62, 1.0, 0.33)]
        for name, bx, by, bw, bh in boxes:
            add_box(slide, x + w * bx + Cm(0.15), y + h * by + Cm(0.15), w * bw - Cm(0.25), h * bh - Cm(0.25), name, fill=PALE_BLUE if name == "전장환경" else WHITE, line=BLUE, size=7, bold=True)
        add_mini_label(slide, "선택 객체 동기화", x + Cm(0.45), y + h - Cm(0.78), Cm(3.1), fill=WHITE)
    elif kind == "underwater":
        for i, label in enumerate(["수온 레이어", "전달손실", "음선경로"]):
            add_box(slide, x + Cm(0.55 + i * 2.45), y + Cm(0.75 + i * 0.5), Cm(2.15), Cm(0.5), label, fill=PALE_BLUE if i % 2 == 0 else WHITE, line=BLUE, size=6.1, bold=True)
        conn = slide.shapes.add_connector(1, int(x + Cm(0.8)), int(y + h - Cm(1.0)), int(x + w - Cm(0.8)), int(y + Cm(1.3)))
        conn.line.color.rgb = BLUE; conn.line.width = Pt(2)
        add_mini_label(slide, "어뢰/기만 분석 연계", x + w - Cm(4.1), y + h - Cm(0.78), Cm(3.3), fill=WHITE)
    elif kind == "model":
        add_box(slide, x + Cm(0.45), y + Cm(0.65), Cm(3.0), Cm(0.72), "입력 파라미터", fill=PALE_BLUE, line=BLUE, size=6.5, bold=True)
        add_box(slide, x + Cm(4.0), y + Cm(0.65), Cm(3.0), Cm(0.72), "모델/산출식", fill=WHITE, line=BLUE, size=6.5, bold=True)
        add_box(slide, x + Cm(7.55), y + Cm(0.65), Cm(3.0), Cm(0.72), "중간 산출값", fill=WHITE, line=BLUE, size=6.5, bold=True)
        add_box(slide, x + Cm(11.1), y + Cm(0.65), Cm(3.0), Cm(0.72), "결과/근거", fill=PALE_BLUE, line=BLUE, size=6.5, bold=True)
        add_mini_label(slide, "입력-산출 추적", x + Cm(5.8), y + Cm(2.45), Cm(3.0), fill=WHITE)
    elif kind in ("chart", "timeline"):
        left, top, cw, ch = x + Cm(0.45), y + Cm(0.85), w - Cm(0.9), h - Cm(1.5)
        for i in range(4):
            grid = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, int(left), int(top + ch * i / 3), int(cw), int(Cm(0.012)))
            grid.fill.solid(); grid.fill.fore_color.rgb = LINE; grid.line.color.rgb = LINE
        last = None
        for px, py in [(0, .7), (.2, .55), (.4, .62), (.6, .38), (.8, .3), (1, .42)]:
            sx, sy = left + cw * px, top + ch * py
            dot = slide.shapes.add_shape(MSO_SHAPE.OVAL, int(sx - Cm(.05)), int(sy - Cm(.05)), int(Cm(.1)), int(Cm(.1)))
            dot.fill.solid(); dot.fill.fore_color.rgb = BLUE; dot.line.color.rgb = BLUE
            if last:
                c = slide.shapes.add_connector(1, int(last[0]), int(last[1]), int(sx), int(sy))
                c.line.color.rgb = BLUE; c.line.width = Pt(1.4)
            last = (sx, sy)
        add_mini_label(slide, "이벤트", x + Cm(0.7), y + h - Cm(0.78), Cm(1.6), fill=WHITE)
        add_mini_label(slide, "임계값", x + Cm(2.55), y + h - Cm(0.78), Cm(1.8), fill=WHITE)
        add_mini_label(slide, "3D 동기화", x + Cm(4.65), y + h - Cm(0.78), Cm(2.2), fill=WHITE)


def add_core_requirement_slide(slide, plan: SlidePlan, width, height):
    detail = CORE_SLIDE_DETAILS[plan.no]
    add_section_label(slide, "구현 개념", Cm(0.95), Cm(4.75))
    add_form_process(slide, detail["steps"], width, y=Cm(5.55))
    add_section_label(slide, "제안 범위 및 검증 관점", Cm(0.95), Cm(8.25), w=Cm(5.0))
    add_form_table(
        slide,
        ["구분", "내용"],
        detail["table"],
        Cm(0.95),
        Cm(8.65),
        [Cm(3.1), width - Cm(5.0)],
        row_h=Cm(0.95),
    )
    add_section_label(slide, "참고 도식", Cm(0.95), Cm(12.9), w=Cm(4.0))
    figure_y = Cm(13.32)
    cols = [
        ("입력자료", "모의 로그\n환경자료\n무기/플랫폼 상태"),
        ("처리기능", "\n".join(detail["steps"][:3])),
        ("분석산출", "3D/차트 전시\n확률·효과 값\n분석 리포트"),
    ]
    box_w = (width - Cm(2.5)) / 3
    for idx, (title, body) in enumerate(cols):
        x0 = Cm(0.95) + idx * (box_w + Cm(0.3))
        add_box(slide, x0, figure_y, box_w, Cm(0.55), title, fill=NAVY if idx == 1 else BLUE, line=BLUE, size=6.5, bold=True, color=WHITE)
        add_box(slide, x0, figure_y + Cm(0.62), box_w, Cm(1.45), body, fill=WHITE if idx != 1 else PALE_BLUE, line=LINE, size=5.4, color=DARK)
        if idx < 2:
            add_arrow(slide, x0 + box_w + Cm(0.05), figure_y + Cm(1.25), Cm(0.22), Cm(0.42))


def create_ppt(images: List[ImageCandidate], plans: List[SlidePlan]) -> None:
    prs = Presentation()
    if FORM_PPTX.exists():
        template = Presentation(FORM_PPTX)
        prs.slide_width = template.slide_width
        prs.slide_height = template.slide_height
    else:
        prs.slide_width = Cm(17.78)
        prs.slide_height = Cm(25.69)
    width, height = prs.slide_width, prs.slide_height
    blank = prs.slide_layouts[6]

    for page_no, plan in enumerate(plans, start=1):
        slide = prs.slides.add_slide(blank)
        if plan.no == 1:
            apply_custom_background(slide, width, height, cover=True)
            t = slide.shapes.add_textbox(Cm(1.1), Cm(6.0), width - Cm(2.2), Cm(2.6))
            set_text(t, "해군교전분석모델 Build-II\n결과분석 제안서", size=22, bold=True, color=DARK)
            sub = slide.shapes.add_textbox(Cm(1.15), Cm(9.1), width - Cm(2.3), Cm(1.0))
            set_text(sub, "RFP 요구사항 대응 중심 | 결과분석 고도화 수행방안", size=10, color=NAVY)
            info = slide.shapes.add_textbox(Cm(1.15), height - Cm(1.55), width - Cm(2.3), Cm(0.55))
            set_text(info, "결과분석 구조 개선 및 구성 기능 고도화 제안", size=8, color=GRAY)
            continue
        if plan.no == 2:
            apply_custom_background(slide, width, height)
            roman = slide.shapes.add_textbox(Cm(1.05), Cm(6.7), Cm(3.0), Cm(2.0))
            set_text(roman, "Ⅱ", size=34, bold=True, color=NAVY, align=PP_ALIGN.CENTER)
            title = slide.shapes.add_textbox(Cm(4.3), Cm(6.95), width - Cm(5.3), Cm(0.9))
            set_text(title, "전략 및 방법론", size=18, bold=True, color=DARK)
            subtitle = slide.shapes.add_textbox(Cm(4.35), Cm(8.05), width - Cm(5.3), Cm(1.0))
            set_text(subtitle, "사업 이해도 / 추진전략 / 결과분석 수행방안 / 시험평가 대응", size=10, color=GRAY)
            line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, int(Cm(4.35)), int(Cm(7.9)), int(width - Cm(5.4)), int(Cm(0.04)))
            line.fill.solid(); line.fill.fore_color.rgb = BLUE; line.line.color.rgb = BLUE
            continue
        add_form_header(slide, plan, width, height)
        add_form_intro(slide, plan, width)
        if plan.no in CORE_SLIDE_DETAILS:
            add_core_requirement_slide(slide, plan, width, height)
        elif plan.no in (4, 24, 28):
            rows = [
                ["3D/자동재생", "로그 기반 3D 재현", "시간축·이벤트 재현성"],
                ["멀티뷰/수중환경", "동기화 뷰·환경 레이어", "입력자료-화면 대조"],
                ["확률/전자전", "실시간 차트·효과 산출", "산출식·시나리오 검증"],
                ["시험평가", "요구사항별 검증 및 산출물 확인", "시험 결과·검증 증적"],
            ]
            if plan.no == 28:
                rows = [
                    ["사업관리", "일정·위험·요구사항 관리", "PM/PL 역할"],
                    ["결과분석/모델", "확률·효과 산출 및 모델 검증", "분석모델 담당"],
                    ["3D 가시화/UI", "3D 재생·멀티뷰·차트 구현", "가시화 담당"],
                    ["시험평가", "시험항목·증적·결함 관리", "검증 담당"],
                ]
            add_form_table(slide, ["구분", "대응 방안", "검증/역할"], rows, Cm(0.95), Cm(4.85), [Cm(4.2), Cm(7.0), width - Cm(12.15)], row_h=Cm(1.12))
        elif plan.no in (6, 7, 18, 19):
            add_form_architecture(slide, width) if plan.no == 19 else add_form_process(slide, ["모의 로그", "분석모델", "3D/차트", "검증 리포트"], width, y=Cm(5.45))
        elif plan.no in (8, 9, 10, 11, 21, 22, 23, 25, 26, 27, 29, 30):
            labels = {
                8: ["전장상황", "플랫폼", "무기체계", "분석데이터"],
                9: ["로그 수집", "시간 정렬", "상태 복원", "3D 재생"],
                10: ["전장환경", "함정", "항공기/잠수함", "무기체계"],
                11: ["수온", "분포", "전달손실", "음선경로"],
                21: ["분석", "설계", "구현", "통합시험"],
                22: ["백로그", "스프린트", "검토", "개선"],
                23: ["형상관리", "CI", "자동시험", "배포관리"],
                25: ["요구사항", "모델 검증", "재현성", "증적"],
                26: ["기준", "리뷰", "결함", "변경관리"],
                27: ["착수", "설계", "구현", "시험평가"],
                29: ["해석성", "재현성", "검증성", "확장성"],
                30: ["요구사항", "구현", "검증", "운영"],
            }.get(plan.no, ["입력", "처리", "분석", "출력"])
            add_form_timeline(slide, width) if plan.no == 27 else add_form_process(slide, labels, width, y=Cm(5.45))
        elif plan.no in (12, 15, 16, 17):
            labels = {
                12: ["입력값", "계산단계", "중간값", "출력값"],
                15: ["기만 이벤트", "유도탄 상태", "효과 산출", "결과 차트"],
                16: ["수중환경", "어뢰 상태", "기만기 효과", "결과 비교"],
                17: ["전자전 이벤트", "재머", "GPS 교란", "효과 분석"],
            }[plan.no]
            add_form_process(slide, labels, width, y=Cm(5.45))
        elif plan.no in (13, 14):
            add_form_chart(slide, width, "확률 산출 결과 실시간 차트")
        else:
            add_form_cards(slide, plan.bullets[:4], width)
        add_form_bullets(slide, plan.bullets, width, height)
        add_form_footer(slide, page_no, width, height)

    try:
        prs.save(PPTX_OUT)
    except PermissionError:
        fallback = PPTX_OUT.with_name(PPTX_OUT.stem + "_정리본.pptx")
        prs.save(fallback)


def main() -> None:
    parser = argparse.ArgumentParser()
    parser.add_argument("--plan-only", action="store_true")
    args = parser.parse_args()

    ensure_dirs()
    task_text = read_task()
    images = render_pdf_assets()
    plans = slide_plan()
    write_markdown(task_text, images, plans)
    if not args.plan_only:
        create_ppt(images, plans)
        print(f"created: {PPTX_OUT}")
    else:
        print("created plan and analysis markdown files")


if __name__ == "__main__":
    main()
