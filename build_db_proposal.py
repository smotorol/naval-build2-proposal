# -*- coding: utf-8 -*-
from __future__ import annotations

from pathlib import Path
from typing import Sequence, Tuple

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_CONNECTOR, MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, MSO_AUTO_SIZE, PP_ALIGN
from pptx.util import Cm, Pt


ROOT = Path(__file__).resolve().parent
FORM_PPTX = ROOT / "form" / "정성제안서.pptx"
TASK_TXT = ROOT / "task.txt"
OUT = ROOT / "해군교전2_기본DB_제안서_고도화.pptx"

NAVY = RGBColor(0x00, 0x4B, 0x73)
BLUE = RGBColor(0x00, 0x79, 0xA8)
SKY = RGBColor(0xD8, 0xEB, 0xF4)
PALE = RGBColor(0xF4, 0xF8, 0xFA)
GRAY = RGBColor(0x66, 0x66, 0x66)
DARK = RGBColor(0x23, 0x2A, 0x31)
LINE = RGBColor(0xC8, 0xD2, 0xDA)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
GREEN = RGBColor(0x69, 0x9B, 0x7A)
RED = RGBColor(0xB8, 0x5C, 0x5C)


def slide_size() -> Tuple[int, int]:
    return Cm(33.867), Cm(19.05)


def read_task() -> str:
    raw = TASK_TXT.read_bytes()
    for enc in ("utf-8", "cp949", "euc-kr"):
        try:
            return raw.decode(enc)
        except UnicodeDecodeError:
            pass
    return raw.decode("utf-8", errors="replace")


def set_text(shape, text: str, size=8, bold=False, color=DARK, align=None) -> None:
    tf = shape.text_frame
    tf.clear()
    tf.word_wrap = True
    tf.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
    tf.margin_left = Cm(0.12)
    tf.margin_right = Cm(0.12)
    tf.margin_top = Cm(0.06)
    tf.margin_bottom = Cm(0.06)
    p = tf.paragraphs[0]
    p.text = text
    p.font.name = "Malgun Gothic"
    if size <= 5.5:
        display_size = size + 1.4
    elif size <= 7.0:
        display_size = size + 1.2
    elif size <= 9.0:
        display_size = size + 0.9
    else:
        display_size = size
    p.font.size = Pt(display_size)
    p.font.bold = bold
    p.font.color.rgb = color
    p.space_after = Pt(0)
    if align:
        p.alignment = align


def textbox(slide, x, y, w, h, text, size=8, bold=False, color=DARK, align=None):
    sh = slide.shapes.add_textbox(int(x), int(y), int(w), int(h))
    set_text(sh, text, size, bold, color, align)
    return sh


def rect(slide, x, y, w, h, fill=WHITE, line=LINE):
    sh = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, int(x), int(y), int(w), int(h))
    sh.fill.solid()
    sh.fill.fore_color.rgb = fill
    sh.line.color.rgb = line
    return sh


def box(slide, x, y, w, h, text="", fill=WHITE, line=LINE, size=7, bold=False, color=DARK, align=PP_ALIGN.CENTER):
    sh = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, int(x), int(y), int(w), int(h))
    sh.fill.solid()
    sh.fill.fore_color.rgb = fill
    sh.line.color.rgb = line
    set_text(sh, text, size, bold, color, align)
    sh.text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
    return sh


def line(slide, x1, y1, x2, y2, color=BLUE, width=1.2):
    c = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, int(x1), int(y1), int(x2), int(y2))
    c.line.color.rgb = color
    c.line.width = Pt(width)
    return c


def arrow(slide, x, y, w, h, color=BLUE):
    sh = slide.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, int(x), int(y), int(w), int(h))
    sh.fill.solid()
    sh.fill.fore_color.rgb = color
    sh.line.color.rgb = color
    return sh


def add_header(slide, num: int, title: str, message: str, width, height):
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = WHITE
    textbox(slide, Cm(0.7), Cm(0.55), Cm(2.2), Cm(0.35), f"Ⅱ - {num-1}", 7.5, True, NAVY)
    textbox(slide, Cm(0.7), Cm(0.95), width - Cm(1.4), Cm(0.62), title, 12.5, True, DARK)
    rect(slide, Cm(0.7), Cm(1.78), width - Cm(1.4), Cm(0.03), NAVY, NAVY)
    rect(slide, Cm(0.7), Cm(1.84), Cm(4.3), Cm(0.07), BLUE, BLUE)
    textbox(slide, Cm(0.9), Cm(2.12), width - Cm(1.8), Cm(0.56), message, 7.8, True, NAVY)


def add_footer(slide, num: int, width, height):
    textbox(slide, Cm(0.7), height - Cm(0.62), width - Cm(1.4), Cm(0.26),
            f"해군교전분석모델 Build-II 기본DB 제안서    {num:02d}", 5.8, False, GRAY, PP_ALIGN.RIGHT)


def process(slide, labels: Sequence[str], x, y, w, h):
    gap = Cm(0.28)
    bw = (w - gap * (len(labels) - 1)) / len(labels)
    for i, label in enumerate(labels):
        xx = x + i * (bw + gap)
        box(slide, xx, y, bw, h, label, fill=WHITE, line=BLUE, size=6.8, bold=True, color=NAVY)
        if i < len(labels) - 1:
            arrow(slide, xx + bw + Cm(0.04), y + h / 2 - Cm(0.14), Cm(0.18), Cm(0.28))


def table(slide, headers: Sequence[str], rows: Sequence[Sequence[str]], x, y, widths: Sequence, row_h=Cm(0.78), size=6.0):
    for i, header in enumerate(headers):
        box(slide, x + sum(widths[:i]), y, widths[i], row_h, header, fill=NAVY, line=NAVY, size=size, bold=True, color=WHITE)
    for r, row in enumerate(rows):
        for c, val in enumerate(row):
            cell = box(slide, x + sum(widths[:c]), y + row_h * (r + 1), widths[c], row_h, val,
                       fill=PALE if r % 2 == 0 else WHITE, line=LINE, size=size - 0.3, color=DARK)
            if c > 0:
                cell.text_frame.paragraphs[0].alignment = PP_ALIGN.LEFT
                cell.text_frame.margin_left = Cm(0.14)


def bullet_panel(slide, bullets: Sequence[str], width, height):
    y = height - Cm(4.05)
    textbox(slide, Cm(0.75), y, Cm(3.0), Cm(0.35), "핵심 포인트", 8, True, NAVY)
    sh = rect(slide, Cm(0.7), y + Cm(0.44), width - Cm(1.4), Cm(2.55), PALE, LINE)
    tf = sh.text_frame
    tf.clear()
    tf.word_wrap = True
    tf.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
    tf.margin_left = Cm(0.28)
    tf.margin_right = Cm(0.28)
    tf.margin_top = Cm(0.12)
    for i, b in enumerate(bullets[:4]):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = f"• {b}"
        p.font.name = "Malgun Gothic"
        p.font.size = Pt(8.2)
        p.font.color.rgb = DARK
        p.alignment = PP_ALIGN.LEFT
        p.space_after = Pt(1)


def title_slide(slide, width, height):
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = WHITE
    rect(slide, Cm(0.85), Cm(5.35), width - Cm(1.7), Cm(0.08), BLUE, BLUE)
    textbox(slide, Cm(1.0), Cm(5.95), width - Cm(2), Cm(2.0), "기본DB\n체계 개선 제안서", 24, True, DARK)
    textbox(slide, Cm(1.05), Cm(8.95), width - Cm(2.1), Cm(0.9),
            "해군교전분석모델 Build-II | 환경·지형정보 및 제원·파라미터 DB 관리 개선", 9.6, False, NAVY)
    process(slide, ["요구사항 정제", "DB 구조 설계", "UI/편집 개선", "검증/동기화", "운영관리"], Cm(1.05), Cm(12.1), width - Cm(2.1), Cm(1.15))
    textbox(slide, Cm(1.05), height - Cm(1.6), width - Cm(2.1), Cm(0.4),
            "기본DB의 표준화, 최신화, 검증 가능성을 확보하는 체계 구축", 8.2, False, GRAY)


def section_slide(slide, width, height):
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = WHITE
    textbox(slide, Cm(1.0), Cm(6.5), Cm(2.8), Cm(1.6), "Ⅱ", 34, True, NAVY, PP_ALIGN.CENTER)
    textbox(slide, Cm(4.2), Cm(6.85), width - Cm(5.2), Cm(0.9), "기본DB 개선 전략", 18, True, DARK)
    rect(slide, Cm(4.25), Cm(7.85), width - Cm(5.4), Cm(0.05), BLUE, BLUE)
    textbox(slide, Cm(4.25), Cm(8.25), width - Cm(5.4), Cm(1.1),
            "환경/지형정보 · 제원/파라미터 · DB 검증 · 사용자관리 · 최신화", 10, False, GRAY)


def db_tobe(slide, width):
    process(slide, ["원천 DB/요구", "표준 스키마", "편집 UI", "검증 규칙", "운영 DB"], Cm(0.9), Cm(4.25), width - Cm(1.8), Cm(1.1))
    layers = [
        ("환경/지형정보", ["공중/해양정보", "작전/임무구역", "격자/해역 속성"]),
        ("제원/파라미터", ["부대/플랫폼", "추진·센서", "전자전·무기체계"]),
        ("품질/운영", ["DB 검증", "권한/이력", "최신화/동기화"]),
    ]
    y = Cm(6.4)
    for name, items in layers:
        box(slide, Cm(1.0), y, Cm(3.5), Cm(0.88), name, fill=NAVY, line=NAVY, size=6.8, bold=True, color=WHITE)
        item_w = (width - Cm(6.0)) / len(items)
        for i, item in enumerate(items):
            box(slide, Cm(4.8) + i * item_w, y, item_w - Cm(0.18), Cm(0.88), item, fill=PALE if i % 2 == 0 else WHITE, line=LINE, size=6.4, bold=True)
        y += Cm(1.3)


def db_schema(slide, width):
    tables = {
        "ENV_GRID": ["grid_id", "sea_area", "depth", "source_ver"],
        "MISSION_AREA": ["area_id", "type", "boundary", "owner"],
        "PLATFORM_SPEC": ["platform_id", "class", "rcs", "safe_depth"],
        "PARAM_LOOKUP": ["param_id", "category", "value", "valid_rule"],
    }
    positions = [(Cm(1), Cm(4.2)), (Cm(9.2), Cm(4.2)), (Cm(1), Cm(8.3)), (Cm(9.2), Cm(8.3))]
    for (name, fields), (x, y) in zip(tables.items(), positions):
        box(slide, x, y, Cm(6.5), Cm(0.62), name, fill=NAVY, line=NAVY, size=6.8, bold=True, color=WHITE)
        box(slide, x, y + Cm(0.72), Cm(6.5), Cm(2.25), "\n".join(fields), fill=PALE, line=LINE, size=6.2, align=PP_ALIGN.LEFT)
    line(slide, Cm(7.5), Cm(5.75), Cm(9.2), Cm(5.75), BLUE, 1.2)
    line(slide, Cm(7.5), Cm(9.8), Cm(9.2), Cm(9.8), BLUE, 1.2)
    line(slide, Cm(4.25), Cm(7.1), Cm(4.25), Cm(8.3), BLUE, 1.2)


def ui_mock(slide, title, tabs: Sequence[str], rows: Sequence[str], x, y, w, h):
    rect(slide, x, y, w, h, WHITE, LINE)
    textbox(slide, x + Cm(0.25), y + Cm(0.18), w - Cm(0.5), Cm(0.35), title, 6.8, True, NAVY, PP_ALIGN.LEFT)
    tab_w = (w - Cm(0.7)) / len(tabs)
    for i, tab in enumerate(tabs):
        box(slide, x + Cm(0.35) + i * tab_w, y + Cm(0.68), tab_w - Cm(0.08), Cm(0.45), tab,
            fill=NAVY if i == 0 else PALE, line=BLUE, size=5.5, bold=True, color=WHITE if i == 0 else NAVY)
    top = y + Cm(1.35)
    widths = [w * 0.25, w * 0.25, w * 0.25, w * 0.20]
    table(slide, ["항목", "값", "검증", "이력"], [[r, "입력/선택", "정상", "기록"] for r in rows[:4]],
          x + Cm(0.35), top, widths, Cm(0.5), size=5.0)


def quality_flow(slide, width):
    process(slide, ["입력", "형식 검증", "참조 무결성", "업무 규칙", "승인/반영"], Cm(0.9), Cm(4.35), width - Cm(1.8), Cm(1.05))
    table(slide, ["검증 관점", "점검 내용", "결과 처리"], [
        ["형식", "필수값, 코드값, 단위, 범위", "오류 표시 및 저장 제한"],
        ["관계", "플랫폼-센서-무장 참조 무결성", "누락 참조 보정"],
        ["업무", "최소/최대값, 운용 조건, 시뮬레이션 사용성", "승인 전 검토"],
        ["이력", "변경자, 변경사유, 적용일자", "버전별 추적"],
    ], Cm(1.0), Cm(6.55), [Cm(3.3), Cm(7.1), width - Cm(12.4)], Cm(0.8), size=5.8)


def body(slide, kind: str, width, height):
    if kind == "understanding":
        table(slide, ["구분", "요구 이해", "제안 방향"], [
            ["환경/지형정보", "해군 통합해양정보체계 DB 구조 연계 및 최신화 용이성", "동일 구조 기반 수용, 갱신 절차 표준화"],
            ["제원/파라미터", "플랫폼·센서·무장·전자전 등 항목 확장", "카테고리별 표준 입력/검증 체계 구성"],
            ["검증/운영", "DB 검증, 사용자관리, 편집 기능 개선", "권한·이력·승인 중심 운영관리"],
        ], Cm(0.9), Cm(4.2), [Cm(3.1), Cm(6.9), width - Cm(11.8)], Cm(1.0), size=5.8)
        ui_mock(slide, "기본DB 화면 예시", ["환경정보", "제원", "검증", "이력"], ["해역", "플랫폼", "센서", "무장"], Cm(1.0), Cm(8.8), width - Cm(2.0), Cm(5.0))
    elif kind == "scope":
        table(slide, ["업무영역", "주요 개선대상", "산출물"], [
            ["환경/지형정보", "공중/해양정보, 작전/임무구역", "표준 테이블, 구역 편집 UI"],
            ["제원/파라미터", "부대, 플랫폼, 추진체계, 센서, 링크", "항목 정의서, 입력/검증 화면"],
            ["무기/전자전", "전자전장비, 화력무기체계, 지휘무장통제", "룩업테이블, 파라미터 관리"],
            ["운영관리", "DB 검증, 사용자 권한, 최신화", "검증규칙, 변경이력, 승인절차"],
        ], Cm(0.9), Cm(4.25), [Cm(3.1), Cm(6.4), width - Cm(11.3)], Cm(0.9), size=5.8)
        process(slide, ["요구분류", "데이터 모델", "화면/편집", "검증", "운영전환"], Cm(1.0), Cm(9.8), width - Cm(2.0), Cm(1.05))
    elif kind == "strategy":
        db_tobe(slide, width)
    elif kind == "schema":
        db_schema(slide, width)
    elif kind == "standard":
        table(slide, ["관리 원칙", "적용 방안"], [
            ["표준 구조", "공통 코드, 식별자, 단위, 좌표계, 버전 정보를 표준화"],
            ["최신화 용이성", "원천 DB 구조와 매핑 규칙을 관리하여 갱신 부담 최소화"],
            ["추적성", "입력자, 변경일자, 변경사유, 적용 버전 이력을 보존"],
            ["검증성", "저장 전 형식·범위·참조·업무규칙 검증을 적용"],
        ], Cm(1.0), Cm(4.3), [Cm(4.0), width - Cm(6.0)], Cm(0.9), size=6.0)
        process(slide, ["코드/단위 표준", "매핑", "검증", "승인", "배포"], Cm(1.0), Cm(9.2), width - Cm(2.0), Cm(1.05))
    elif kind == "env":
        table(slide, ["항목", "개선 내용", "관리 방향"], [
            ["공중/해양정보", "통합해양정보체계 DB와 동일 구조 사용", "해역별 격자/수심/환경속성 표준화"],
            ["최신화", "향후 DB 최신화 용이 목적", "원천-운영 DB 매핑 및 이력 관리"],
            ["조위/환경", "조위 정보 정렬·평균 기준 변환 기능", "자동 산출값과 원천값 병행 관리"],
        ], Cm(0.9), Cm(4.2), [Cm(3.1), Cm(6.7), width - Cm(11.6)], Cm(0.9), size=5.8)
        ui_mock(slide, "환경/지형정보 관리 예시", ["해역", "격자", "수심", "조위"], ["해역코드", "격자수준", "수심", "원천버전"], Cm(1.0), Cm(8.6), width - Cm(2.0), Cm(5.0))
    elif kind == "area":
        table(slide, ["구분", "현재 개선 필요", "제안 방향"], [
            ["작전/임무구역", "구역설정 UI 개선 필요", "참조구역/선 관리와 동일한 편집 패턴 적용"],
            ["편집 기능", "구역 경계, 선, 속성 관리 일관성 부족", "그리기·수정·저장·검증 절차 통합"],
            ["운영성", "다수 구역 변경 시 이력 추적 필요", "변경 이력 및 승인 상태 관리"],
        ], Cm(1.0), Cm(4.25), [Cm(3.4), Cm(5.7), width - Cm(11.1)], Cm(0.9), size=5.8)
        ui_mock(slide, "구역 설정 UI 개선 예시", ["참조구역", "선 관리", "임무구역"], ["구역명", "좌표", "임무유형", "상태"], Cm(1.0), Cm(8.8), width - Cm(2.0), Cm(5.0))
    elif kind == "params":
        table(slide, ["분류", "추가/개선 대상", "관리 방식"], [
            ["부대", "제3국 부대 추가", "수상함/항공/잠수함 부대 유형화"],
            ["플랫폼", "Wave glider, 저궤도 위성 추가", "플랫폼별 속성/형상/RCS 관리"],
            ["추진체계", "무인체계 배터리·엔진, 무인기 배터리", "동력원 및 운용시간 파라미터화"],
            ["센서", "SAR 추가", "센서 유형/성능/제한각 관리"],
        ], Cm(0.9), Cm(4.2), [Cm(2.7), Cm(6.3), width - Cm(10.8)], Cm(0.82), size=5.7)
        process(slide, ["분류체계", "속성 정의", "룩업 관리", "검증 규칙", "적용 이력"], Cm(1.0), Cm(9.4), width - Cm(2.0), Cm(1.05))
    elif kind == "platform":
        table(slide, ["항목", "제안 내용"], [
            ["플랫폼 확장", "Wave glider, 저궤도 위성 등 신규 플랫폼 유형 추가"],
            ["3D 형상/RCS", "수상함·잠수함·항공기 형상 기반 RCS 산출 항목 관리"],
            ["안전항해수심", "수상항해/스노클항해/잠항시 조건별 수심 속성 관리"],
            ["플랫폼 속성", "진영, 운용상태, 탑재체계, 시스템 네트워크 속성 연계"],
        ], Cm(1.0), Cm(4.25), [Cm(3.6), width - Cm(5.6)], Cm(0.82), size=5.8)
        ui_mock(slide, "플랫폼 제원 편집 예시", ["기본", "형상", "RCS", "운용"], ["플랫폼ID", "형상모델", "대표RCS", "안전수심"], Cm(1.0), Cm(8.6), width - Cm(2.0), Cm(5.0))
    elif kind == "propulsion":
        table(slide, ["관리대상", "개선 내용", "검증 관점"], [
            ["추진체계", "배터리, 엔진 등 무인체계 추진체계 추가", "운용시간/동력원 범위 검증"],
            ["항공기", "무인기 배터리 속성 반영", "임무 지속시간 연계"],
            ["데이터링크", "전송지연, 위성 up/down-link 지연, 손실률", "지연/정확도 단위 및 범위 검증"],
        ], Cm(0.9), Cm(4.25), [Cm(3.2), Cm(6.7), width - Cm(11.7)], Cm(0.9), size=5.7)
        process(slide, ["동력원", "지연시간", "정확도", "손실률", "시뮬레이션 반영"], Cm(1.0), Cm(9.5), width - Cm(2.0), Cm(1.05))
    elif kind == "ew":
        table(slide, ["구분", "DB 관리 항목", "제안 방향"], [
            ["EA", "레이더/통신체계 교란, 재밍효과", "방식별 제원값 또는 모델결과 룩업테이블화"],
            ["EP", "전자공격 저항 방호능력", "센서/통신체계별 생존성 평가 파라미터 관리"],
            ["검증", "교전급 수준 제원값, 참조 테이블", "입력값 범위와 참조 무결성 검증"],
        ], Cm(0.9), Cm(4.25), [Cm(2.8), Cm(6.6), width - Cm(11.2)], Cm(0.9), size=5.7)
        ui_mock(slide, "전자전장비 룩업테이블 예시", ["EA", "EP", "룩업", "검증"], ["방식", "대상체계", "효과값", "적용조건"], Cm(1.0), Cm(8.75), width - Cm(2.0), Cm(4.9))
    elif kind == "weapon":
        table(slide, ["분류", "추가/개선 대상", "관리 방향"], [
            ["화력무기체계", "항적추적어뢰, 초공동어뢰, 설정심도", "무기별 사거리·심도·센서부 속성 관리"],
            ["지휘무장통제", "다수 TG 운용, DDG I/II 동시교전능력", "교전능력 구분 및 복합전 지휘구조 모델링"],
            ["무장통제", "대탄도탄/유도탄 교전 개념 반영", "동시교전 제한 및 할당 기준 데이터화"],
        ], Cm(0.9), Cm(4.25), [Cm(3.2), Cm(6.7), width - Cm(11.7)], Cm(0.9), size=5.7)
        process(slide, ["무기분류", "운용조건", "교전제한", "능력값", "모의 반영"], Cm(1.0), Cm(9.6), width - Cm(2.0), Cm(1.05))
    elif kind == "damage":
        table(slide, ["항목", "개선 내용", "제안 방향"], [
            ["표준 플랫폼", "함형/톤수/무기체계별 표준 플랫폼 3D 및 네트워크 모듈화", "플랫폼 기준 손상평가 DB 관리"],
            ["입력 방식", "컬럼 입력 중심으로 3D 직관성 부족", "입력값-3D 표현 연계 UI 구성"],
            ["진영 변경", "플랫폼 진영 변경 시 손상평가 자동 변경", "진영 속성과 평가 기준 연동"],
            ["신규 플랫폼", "적 및 신규 플랫폼 손상평가용 3D 렌더링/네트워크 자동생성", "템플릿 기반 자동 생성"],
        ], Cm(0.9), Cm(4.2), [Cm(3.2), Cm(6.7), width - Cm(11.7)], Cm(0.82), size=5.5)
        ui_mock(slide, "손상평가 DB 편집 예시", ["플랫폼", "구조", "평가기준", "3D"], ["함형", "톤수", "무기체계", "진영"], Cm(1.0), Cm(8.9), width - Cm(2.0), Cm(4.7))
    elif kind == "validation":
        quality_flow(slide, width)
    elif kind == "user":
        table(slide, ["관리 항목", "제안 내용"], [
            ["권한", "다수 사용자에 대해 일괄 권한 조정 기능 제공"],
            ["삭제/승인", "삭제 및 반영 권한을 역할별로 통제"],
            ["편집성", "테이블ID, 메뉴ID 가독성 및 편집 기능 개선"],
            ["감사이력", "변경자, 변경 전/후 값, 승인자, 적용일자 추적"],
        ], Cm(1.0), Cm(4.3), [Cm(3.6), width - Cm(5.6)], Cm(0.85), size=5.8)
        process(slide, ["사용자", "역할", "권한", "승인", "감사이력"], Cm(1.0), Cm(9.4), width - Cm(2.0), Cm(1.05))
    elif kind == "migration":
        process(slide, ["원천 확인", "매핑", "검증", "승인", "운영 반영"], Cm(0.9), Cm(4.3), width - Cm(1.8), Cm(1.05))
        table(slide, ["단계", "주요 내용", "산출물"], [
            ["원천 확인", "통합해양정보체계 및 기존 기본DB 구조 확인", "원천/대상 매핑표"],
            ["매핑", "테이블·컬럼·코드·단위 변환 규칙 정의", "변환 규칙서"],
            ["검증", "건수/참조/범위/업무규칙 검증", "검증 결과"],
            ["반영", "승인 후 운영 DB 반영 및 이력 저장", "반영 이력"],
        ], Cm(1.0), Cm(6.6), [Cm(2.7), Cm(7.0), width - Cm(11.7)], Cm(0.8), size=5.8)
    elif kind == "ui":
        ui_mock(slide, "기본DB 통합관리 UI 방향", ["검색", "편집", "검증", "이력"], ["분류", "속성", "검증결과", "승인상태"], Cm(1.0), Cm(4.25), width - Cm(2.0), Cm(6.2))
        table(slide, ["UI 원칙", "적용 방향"], [
            ["동일 패턴", "참조구역/선 관리와 유사한 편집 경험 제공"],
            ["오류 가시화", "입력 즉시 오류 위치와 사유 표시"],
            ["변경 추적", "변경 전/후 값과 승인 상태를 한 화면에서 확인"],
        ], Cm(1.0), Cm(11.1), [Cm(3.6), width - Cm(5.6)], Cm(0.78), size=5.8)
    elif kind == "plan":
        table(slide, ["단계", "수행 내용", "산출물"], [
            ["분석", "task 및 RFP 기본DB 요구사항 정제", "요구사항 매트릭스"],
            ["설계", "DB 스키마, 코드, UI, 검증규칙 설계", "DB/UI 설계서"],
            ["구현", "환경/제원/무기/전자전/손상평가 관리 화면 구현", "기능 모듈"],
            ["검증", "데이터 이관, 정합성, 권한, 이력 시험", "검증 결과서"],
            ["전환", "운영 DB 반영 및 사용자 교육", "전환 계획서"],
        ], Cm(0.9), Cm(4.2), [Cm(2.4), Cm(7.5), width - Cm(11.7)], Cm(0.78), size=5.6)
        process(slide, ["분석", "설계", "구현", "검증", "전환"], Cm(1.0), Cm(9.6), width - Cm(2.0), Cm(1.05))
    elif kind == "risk":
        table(slide, ["위험요소", "영향", "대응방안"], [
            ["원천 DB 구조 차이", "최신화 지연", "매핑 규칙과 예외 처리 기준 수립"],
            ["입력항목 증가", "사용성 저하", "분류별 화면 분리와 일괄 편집 제공"],
            ["참조 무결성 오류", "모의 실행 오류", "저장 전 검증 및 승인 절차 적용"],
            ["권한관리 미흡", "오입력/오삭제 위험", "역할 기반 권한과 감사이력 적용"],
        ], Cm(0.9), Cm(4.3), [Cm(3.6), Cm(4.8), width - Cm(10.2)], Cm(0.9), size=5.6)
        process(slide, ["식별", "분석", "대응", "검증", "이력화"], Cm(1.0), Cm(9.8), width - Cm(2.0), Cm(1.05))
    elif kind == "effect":
        table(slide, ["기대효과", "설명"], [
            ["최신화 용이성", "원천 DB 구조와 매핑 규칙 기반으로 갱신 부담 완화"],
            ["입력 품질 향상", "형식·범위·참조 검증으로 오입력 최소화"],
            ["운영 추적성", "변경자, 승인자, 버전, 변경사유 이력 확보"],
            ["확장성", "신규 플랫폼, 센서, 무장, 전자전 항목 추가에 대응"],
        ], Cm(1.0), Cm(4.4), [Cm(3.5), width - Cm(5.5)], Cm(0.9), size=6.0)
        process(slide, ["표준화", "검증", "이력", "확장", "운영 안정화"], Cm(1.0), Cm(9.4), width - Cm(2.0), Cm(1.05))
    elif kind == "conclusion":
        table(slide, ["구분", "최종 제안 메시지"], [
            ["기본DB 표준화", "환경/지형정보와 제원/파라미터를 표준 구조로 관리"],
            ["최신화 체계", "통합해양정보체계 등 원천 구조와 연계 가능한 갱신 기반 확보"],
            ["검증 중심 운영", "저장 전 검증, 권한, 이력, 승인 절차로 운영 안정성 강화"],
            ["확장 기반", "신규 플랫폼·센서·무장·전자전·손상평가 항목 추가에 대응"],
        ], Cm(1.0), Cm(4.7), [Cm(4.1), width - Cm(6.1)], Cm(0.9), size=6.0)
        process(slide, ["요구 충족", "구조 표준화", "품질 검증", "운영관리", "확장성"], Cm(1.0), Cm(9.4), width - Cm(2.0), Cm(1.05))


SLIDES = [
    ("표지", "", "cover"),
    ("사업 이해", "Build-II 기본DB는 환경/지형정보와 제원/파라미터의 표준화 및 최신화 기반 확보가 핵심입니다.", "understanding"),
    ("기본DB 업무 범위", "기본DB 대상과 산출물을 환경정보, 제원정보, 검증/운영관리로 구분합니다.", "scope"),
    ("기본DB To-Be 관리 체계", "원천 DB 연계, 표준 스키마, 편집 UI, 검증 규칙, 운영 DB를 일관된 흐름으로 구성합니다.", "strategy"),
    ("기본DB 논리 구조 개념", "환경/구역/플랫폼/파라미터 정보를 추적 가능한 테이블 구조로 관리합니다.", "schema"),
    ("데이터 표준화 및 최신화 전략", "코드, 단위, 식별자, 버전 정보를 표준화하여 향후 DB 최신화를 용이하게 합니다.", "standard"),
    ("환경/지형정보 관리 개선", "공중/해양정보와 해역별 환경 정보를 원천 구조와 연계 가능한 형태로 관리합니다.", "env"),
    ("작전/임무구역 UI 개선", "구역 설정 UI를 참조구역/선 관리와 동일한 편집 경험으로 개선합니다.", "area"),
    ("제원/파라미터 관리 체계", "부대, 플랫폼, 추진체계, 센서, 링크 등 확장 항목을 분류 기반으로 관리합니다.", "params"),
    ("플랫폼 DB 확장", "Wave glider, 저궤도 위성 등 신규 플랫폼과 형상/RCS/안전항해수심 속성을 반영합니다.", "platform"),
    ("추진체계·센서·데이터링크 관리", "무인체계 추진체계, SAR, 전송지연 및 데이터 손실률을 제원 DB로 관리합니다.", "propulsion"),
    ("전자전장비 DB 관리", "EA/EP 방식별 제원값과 룩업테이블을 통해 전자전 관련 기본DB를 관리합니다.", "ew"),
    ("화력무기체계 및 지휘무장통제 DB", "신규 무기체계와 동시교전능력, 복합전 지휘구조를 데이터화합니다.", "weapon"),
    ("손상평가관리 개선", "플랫폼 3D/시스템 네트워크 모듈화와 손상평가 DB 연계를 강화합니다.", "damage"),
    ("DB 검증 및 품질관리", "형식, 참조 무결성, 업무규칙, 이력 검증을 저장 전후 단계에 적용합니다.", "validation"),
    ("사용자·권한·편집 관리", "다수 사용자 권한 조정, 삭제 통제, 편집 기능, 감사이력 관리를 제안합니다.", "user"),
    ("데이터 이관 및 동기화 방안", "원천 확인, 매핑, 검증, 승인, 운영 반영 절차로 DB 최신화를 수행합니다.", "migration"),
    ("UI/UX 설계 방향", "검색, 편집, 검증, 이력 확인을 한 흐름으로 제공하는 기본DB 통합관리 UI를 구성합니다.", "ui"),
    ("개발 및 시험 계획", "분석, 설계, 구현, 검증, 전환 단계로 기본DB 개선을 추진합니다.", "plan"),
    ("위험요소 및 대응방안", "원천 DB 구조 차이, 입력항목 증가, 참조 무결성, 권한관리 리스크를 사전에 관리합니다.", "risk"),
    ("기대효과", "기본DB 최신화 용이성, 입력 품질, 운영 추적성, 확장성을 확보합니다.", "effect"),
    ("결론", "기본DB 표준화와 검증 중심 운영을 통해 Build-II 체계개발 요구사항에 대응합니다.", "conclusion"),
]


BULLETS = {
    "understanding": ["기본DB 요구사항으로 제안 방향 전환", "환경/지형정보와 제원/파라미터 중심 구성", "최신화·검증·운영 역량 강조"],
    "scope": ["기본DB 담당 범위를 명확화", "카테고리별 산출물과 검증 기준 제시", "운영 단계 변경 이력까지 포함"],
    "strategy": ["원천 DB부터 운영 DB까지 흐름화", "편집 UI와 검증 규칙을 함께 설계", "확장 가능한 분류체계 기반 구성"],
    "schema": ["환경, 구역, 플랫폼, 파라미터 핵심 구조 제안", "식별자와 버전 기반 추적성 확보", "참조 무결성 검증 기반 마련"],
    "standard": ["코드/단위/식별자 표준화", "원천-운영 DB 매핑 규칙화", "검증 후 승인 반영 절차 구성"],
    "env": ["통합해양정보체계 DB 구조 연계", "해역별 환경정보 최신화 용이성 확보", "조위/수심/격자 속성 관리"],
    "area": ["구역 편집 UI 일관성 확보", "참조구역/선 관리와 동일한 조작 패턴", "변경 이력과 승인 상태 관리"],
    "params": ["신규 부대/플랫폼/센서/링크 항목 반영", "제원값 입력과 검증 규칙 연계", "룩업테이블 기반 관리"],
    "platform": ["신규 플랫폼과 형상/RCS 속성 확장", "안전항해수심 등 운용 조건 관리", "플랫폼 속성 변경 이력 확보"],
    "propulsion": ["무인체계 추진체계와 배터리 항목 반영", "데이터링크 지연/손실률 관리", "센서 유형 확장에 대응"],
    "ew": ["EA/EP 방식별 룩업테이블 관리", "전자전 제원값과 모델 결과 연계", "참조 무결성 및 범위 검증"],
    "weapon": ["신규 무기체계와 설정심도 항목 반영", "동시교전능력과 지휘구조 데이터화", "복합전 분석 기반 확보"],
    "damage": ["플랫폼 3D/네트워크 모듈화 연계", "손상평가 기준 DB화", "진영 변경에 따른 평가 기준 연동"],
    "validation": ["형식/관계/업무/이력 검증 체계화", "저장 전 오류 차단", "승인 기반 운영 안정성 확보"],
    "user": ["역할 기반 권한관리", "삭제 및 반영 통제", "감사이력으로 운영 추적성 강화"],
    "migration": ["원천-대상 매핑 절차 명확화", "검증 후 승인 반영", "최신화 반복 업무 부담 완화"],
    "ui": ["검색-편집-검증-이력 흐름 제공", "오류 위치와 사유 가시화", "기존 관리 화면과 유사한 사용성"],
    "plan": ["단계별 산출물 중심 추진", "검증과 전환을 별도 단계로 관리", "운영자 교육과 전환 계획 포함"],
    "risk": ["구조 차이와 무결성 오류 사전 관리", "입력항목 증가에 따른 사용성 보완", "권한 및 이력 관리로 오조작 방지"],
    "effect": ["DB 최신화 용이성 향상", "입력 품질과 검증성 강화", "신규 항목 확장 기반 확보"],
    "conclusion": ["기본DB 표준화 중심 제안", "검증 가능한 운영관리 체계 확보", "Build-II 확장 요구에 단계적 대응"],
}


def build():
    _ = read_task()
    prs = Presentation()
    width, height = slide_size()
    prs.slide_width = width
    prs.slide_height = height
    blank = prs.slide_layouts[6]
    for idx, (title, msg, kind) in enumerate(SLIDES, start=1):
        slide = prs.slides.add_slide(blank)
        if kind == "cover":
            title_slide(slide, width, height)
            continue
        if idx == 2:
            section_slide(slide, width, height)
            continue
        add_header(slide, idx, title, msg, width, height)
        body(slide, kind, width, height)
        bullet_panel(slide, BULLETS.get(kind, []), width, height)
        add_footer(slide, idx, width, height)
    prs.save(OUT)


if __name__ == "__main__":
    build()
